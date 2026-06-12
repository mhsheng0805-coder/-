"""
紡織所業務部門來自民間業務收入明細 - Flask Web App
支援多年度（民國年）
執行: python app.py  |  網址: http://[本機IP]:5001
"""
from flask import (Flask, render_template, request, jsonify, redirect,
                   url_for, send_file, session, g)
import sqlite3, os, secrets, hashlib
from datetime import datetime, timedelta
import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter
import io
try:
    from pptx import Presentation as PptxPresentation
    from pptx.util import Inches, Pt, Cm, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn
    from pptx.util import Pt
    import copy
    _PPTX_OK = True
except ImportError:
    _PPTX_OK = False
from functools import wraps

# ── PostgreSQL / SQLite 雙模式 ─────────────────────────
_DATABASE_URL = os.environ.get('DATABASE_URL', '')
if _DATABASE_URL.startswith('postgres://'):
    _DATABASE_URL = _DATABASE_URL.replace('postgres://', 'postgresql://', 1)
IS_PG = bool(_DATABASE_URL)

if IS_PG:
    import psycopg2, psycopg2.extras

class _PgWrapper:
    """讓 psycopg2 連線的介面與 sqlite3 一致（支援 ?  佔位符 + dict row）"""
    def __init__(self, conn):
        self._conn = conn

    def execute(self, sql, params=()):
        sql = sql.replace('?', '%s')
        cur = self._conn.cursor(cursor_factory=psycopg2.extras.RealDictCursor)
        cur.execute(sql, params if params else None)
        return cur

    def commit(self):  self._conn.commit()
    def close(self):   self._conn.close()
    def __enter__(self): return self
    def __exit__(self, *a): self.close()

app = Flask(__name__)
# secret_key 固定化以確保重啟後 session 仍有效
_SECRET_FILE = os.path.join(os.path.dirname(__file__), '.secret_key')
if os.path.exists(_SECRET_FILE):
    app.secret_key = open(_SECRET_FILE).read().strip()
else:
    app.secret_key = secrets.token_hex(32)
    open(_SECRET_FILE, 'w').write(app.secret_key)

app.config['SESSION_COOKIE_HTTPONLY'] = True
app.config['SESSION_COOKIE_SAMESITE'] = 'Lax'
app.config['PERMANENT_SESSION_LIFETIME'] = timedelta(hours=8)

DB = os.path.join(os.path.dirname(__file__), 'revenue.db')

# 當前民國年（預設）
CURRENT_ROC_YEAR = datetime.now().year - 1911

def hash_pw(pw):
    return hashlib.sha256(pw.encode('utf-8')).hexdigest()

def login_required(f):
    @wraps(f)
    def decorated(*args, **kwargs):
        if not session.get('user'):
            return redirect(url_for('login', next=request.path))
        return f(*args, **kwargs)
    return decorated

# ─────────────────────────────────────────────────────
DEPARTMENTS = ['原料部', '產品部', '檢驗部', '製程部', '雲分部', '產服部']
MONTHS = list(range(1, 13))

INCOME_ITEMS = [
    '其他民間收入',
    '配合款-產發署案收入',
    '配合款-其他民間收入',
    '科專衍生收入',
    '能源署衍生收入',
    '其他專案-來自民間',
    '來自民間收入',
]
INCOME_TOTAL_ITEM = '來自民間收入'
EXPENSE_ITEMS = [
    '人事費',
    '業務費',
    '維護費',
    '旅運費',
    '材料費',
    '產發署案配合款',
    '配合款-其他民間',
    '其他專案-民間',
    '轉回專案設備使用費',
    '計畫衍生支出',
    '其他民間收入支出',
]
UNCLAIMED_ITEMS = ['業務費(未核銷)', '旅運費(未核銷)', '材料費(未核銷)', '維護費(未核銷)']
CONTRACT_STATUSES = ['洽談中', '新增簽約', '已簽約執行中', '完成']

# PPT 匯出用：各部門的收入/支出項目（對應企劃處簡報格式）
DEPT_INCOME_PPT = {
    '原料部': ['其他民間收入(試驗/技術/訓練/其他)', '配合款-產發署案收入', '其他專案-來自民間', '科專衍生收入'],
    '產品部': ['其他民間收入(試驗/技術/訓練/其他)', '配合款-產發署案收入', '配合款-其他民間收入', '其他專案-來自民間'],
    '檢驗部': ['其他民間收入(試驗/技術/訓練/其他)', '配合款-產發署案收入', '其他專案-來自民間', '科專衍生收入'],
    '製程部': ['其他民間收入(試驗/技術/訓練/其他)', '配合款-產發署案收入', '其他專案-來自民間', '科專衍生收入', '能源署衍生收入'],
    '雲分部': ['其他民間收入(試驗/技術/訓練/其他)', '配合款-產發署案收入', '配合款-其他民間收入', '其他專案-來自民間', '科專衍生收入'],
    '產服部': ['其他民間收入(試驗/技術/訓練/其他)', '配合款-產發署案收入', '其他專案-來自民間', '科專衍生收入'],
}
DEPT_EXPENSE_PPT = {dept: [
    '人事費', '業務費', '維護費', '旅運費', '材料費',
    '產發署案配合款', '配合款-其他民間', '其他專案-民間',
    '轉回專案設備使用費',
    '計畫衍生支出',
] for dept in DEPARTMENTS}

# ── 資料庫 ─────────────────────────────────────────────
_SCHEMA_SQLITE = '''
CREATE TABLE IF NOT EXISTS settings (key TEXT PRIMARY KEY, value TEXT);
CREATE TABLE IF NOT EXISTS users (
    id INTEGER PRIMARY KEY AUTOINCREMENT,
    username TEXT NOT NULL UNIQUE, password_hash TEXT NOT NULL,
    display_name TEXT, role TEXT DEFAULT 'user', dept TEXT DEFAULT '',
    disabled INTEGER DEFAULT 0,
    reset_token TEXT, reset_expires DATETIME,
    created_at DATETIME DEFAULT CURRENT_TIMESTAMP
);
CREATE TABLE IF NOT EXISTS revenue (
    id INTEGER PRIMARY KEY AUTOINCREMENT,
    year INTEGER NOT NULL, dept TEXT NOT NULL, month INTEGER NOT NULL,
    item TEXT NOT NULL, amount REAL DEFAULT 0, goal REAL DEFAULT 0,
    updated_at DATETIME DEFAULT CURRENT_TIMESTAMP,
    UNIQUE(year, dept, month, item)
);
CREATE TABLE IF NOT EXISTS contracts (
    id INTEGER PRIMARY KEY AUTOINCREMENT,
    year INTEGER NOT NULL, dept TEXT NOT NULL, month INTEGER NOT NULL,
    client TEXT, project_name TEXT DEFAULT '', amount REAL DEFAULT 0, sign_date TEXT,
    status TEXT DEFAULT '洽談中', group_name TEXT DEFAULT '',
    note TEXT, carry_next INTEGER DEFAULT 0,
    cross_dept INTEGER DEFAULT 0, cross_dept_data TEXT DEFAULT '{}',
    payment_type TEXT DEFAULT '當年',
    installments INTEGER DEFAULT 1,
    installment_data TEXT DEFAULT '[]',
    expected_amount REAL DEFAULT 0, expected_date TEXT DEFAULT '',
    updated_at DATETIME DEFAULT CURRENT_TIMESTAMP
);
CREATE TABLE IF NOT EXISTS unclaimed (
    id INTEGER PRIMARY KEY AUTOINCREMENT,
    year INTEGER NOT NULL, dept TEXT NOT NULL, month INTEGER NOT NULL,
    item TEXT NOT NULL, amount REAL DEFAULT 0, note TEXT,
    updated_at DATETIME DEFAULT CURRENT_TIMESTAMP,
    UNIQUE(year, dept, month, item)
);
CREATE TABLE IF NOT EXISTS locks (
    id INTEGER PRIMARY KEY AUTOINCREMENT,
    year INTEGER NOT NULL, dept TEXT NOT NULL, month INTEGER NOT NULL,
    locked INTEGER DEFAULT 0, locked_by TEXT, locked_at DATETIME,
    unlock_requested INTEGER DEFAULT 0, req_by TEXT, req_at DATETIME, req_reason TEXT,
    unlocked_by TEXT, unlocked_at DATETIME,
    UNIQUE(year, dept, month)
);
CREATE TABLE IF NOT EXISTS audit_log (
    id INTEGER PRIMARY KEY AUTOINCREMENT,
    year INTEGER NOT NULL, dept TEXT NOT NULL, month INTEGER NOT NULL,
    action TEXT NOT NULL, table_name TEXT,
    item TEXT, old_value TEXT, new_value TEXT,
    changed_by TEXT, changed_at DATETIME DEFAULT CURRENT_TIMESTAMP,
    note TEXT
);
CREATE TABLE IF NOT EXISTS annual_goals (
    id INTEGER PRIMARY KEY AUTOINCREMENT,
    year INTEGER NOT NULL, dept TEXT NOT NULL, item TEXT NOT NULL,
    goal REAL DEFAULT 0,
    updated_at DATETIME DEFAULT CURRENT_TIMESTAMP,
    UNIQUE(year, dept, item)
);
CREATE TABLE IF NOT EXISTS dept_groups (
    id INTEGER PRIMARY KEY AUTOINCREMENT,
    dept TEXT NOT NULL,
    group_name TEXT NOT NULL,
    sort_order INTEGER DEFAULT 0,
    UNIQUE(dept, group_name)
);
'''

_SCHEMA_PG = '''
CREATE TABLE IF NOT EXISTS settings (key TEXT PRIMARY KEY, value TEXT);
CREATE TABLE IF NOT EXISTS users (
    id SERIAL PRIMARY KEY,
    username TEXT NOT NULL UNIQUE, password_hash TEXT NOT NULL,
    display_name TEXT, role TEXT DEFAULT 'user', dept TEXT DEFAULT '',
    disabled INTEGER DEFAULT 0,
    reset_token TEXT, reset_expires TIMESTAMP,
    created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
);
CREATE TABLE IF NOT EXISTS revenue (
    id SERIAL PRIMARY KEY,
    year INTEGER NOT NULL, dept TEXT NOT NULL, month INTEGER NOT NULL,
    item TEXT NOT NULL, amount REAL DEFAULT 0, goal REAL DEFAULT 0,
    updated_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
    UNIQUE(year, dept, month, item)
);
CREATE TABLE IF NOT EXISTS contracts (
    id SERIAL PRIMARY KEY,
    year INTEGER NOT NULL, dept TEXT NOT NULL, month INTEGER NOT NULL,
    client TEXT, project_name TEXT DEFAULT '', amount REAL DEFAULT 0, sign_date TEXT,
    status TEXT DEFAULT '洽談中', group_name TEXT DEFAULT '',
    note TEXT, carry_next INTEGER DEFAULT 0,
    cross_dept INTEGER DEFAULT 0, cross_dept_data TEXT DEFAULT '{}',
    payment_type TEXT DEFAULT '當年',
    installments INTEGER DEFAULT 1,
    installment_data TEXT DEFAULT '[]',
    expected_amount REAL DEFAULT 0, expected_date TEXT DEFAULT '',
    updated_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
);
CREATE TABLE IF NOT EXISTS unclaimed (
    id SERIAL PRIMARY KEY,
    year INTEGER NOT NULL, dept TEXT NOT NULL, month INTEGER NOT NULL,
    item TEXT NOT NULL, amount REAL DEFAULT 0, note TEXT,
    updated_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
    UNIQUE(year, dept, month, item)
);
CREATE TABLE IF NOT EXISTS locks (
    id SERIAL PRIMARY KEY,
    year INTEGER NOT NULL, dept TEXT NOT NULL, month INTEGER NOT NULL,
    locked INTEGER DEFAULT 0, locked_by TEXT, locked_at TIMESTAMP,
    unlock_requested INTEGER DEFAULT 0, req_by TEXT, req_at TIMESTAMP, req_reason TEXT,
    unlocked_by TEXT, unlocked_at TIMESTAMP,
    UNIQUE(year, dept, month)
);
CREATE TABLE IF NOT EXISTS audit_log (
    id SERIAL PRIMARY KEY,
    year INTEGER NOT NULL, dept TEXT NOT NULL, month INTEGER NOT NULL,
    action TEXT NOT NULL, table_name TEXT,
    item TEXT, old_value TEXT, new_value TEXT,
    changed_by TEXT, changed_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
    note TEXT
);
CREATE TABLE IF NOT EXISTS annual_goals (
    id SERIAL PRIMARY KEY,
    year INTEGER NOT NULL, dept TEXT NOT NULL, item TEXT NOT NULL,
    goal REAL DEFAULT 0,
    updated_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
    UNIQUE(year, dept, item)
);
CREATE TABLE IF NOT EXISTS dept_groups (
    id SERIAL PRIMARY KEY,
    dept TEXT NOT NULL,
    group_name TEXT NOT NULL,
    sort_order INTEGER DEFAULT 0,
    UNIQUE(dept, group_name)
);
'''

_MIGRATE_USERS = [
    "ALTER TABLE users ADD COLUMN dept TEXT DEFAULT ''",
    "ALTER TABLE users ADD COLUMN disabled INTEGER DEFAULT 0",
]
_MIGRATE_NEW_TABLES_SQLITE = [
    """CREATE TABLE IF NOT EXISTS locks (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        year INTEGER NOT NULL, dept TEXT NOT NULL, month INTEGER NOT NULL,
        locked INTEGER DEFAULT 0, locked_by TEXT, locked_at DATETIME,
        unlock_requested INTEGER DEFAULT 0, req_by TEXT, req_at DATETIME, req_reason TEXT,
        unlocked_by TEXT, unlocked_at DATETIME,
        UNIQUE(year, dept, month))""",
    """CREATE TABLE IF NOT EXISTS audit_log (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        year INTEGER NOT NULL, dept TEXT NOT NULL, month INTEGER NOT NULL,
        action TEXT NOT NULL, table_name TEXT,
        item TEXT, old_value TEXT, new_value TEXT,
        changed_by TEXT, changed_at DATETIME DEFAULT CURRENT_TIMESTAMP, note TEXT)""",
]
_MIGRATE_NEW_TABLES_PG = [
    """CREATE TABLE IF NOT EXISTS locks (
        id SERIAL PRIMARY KEY,
        year INTEGER NOT NULL, dept TEXT NOT NULL, month INTEGER NOT NULL,
        locked INTEGER DEFAULT 0, locked_by TEXT, locked_at TIMESTAMP,
        unlock_requested INTEGER DEFAULT 0, req_by TEXT, req_at TIMESTAMP, req_reason TEXT,
        unlocked_by TEXT, unlocked_at TIMESTAMP,
        UNIQUE(year, dept, month))""",
    """CREATE TABLE IF NOT EXISTS audit_log (
        id SERIAL PRIMARY KEY,
        year INTEGER NOT NULL, dept TEXT NOT NULL, month INTEGER NOT NULL,
        action TEXT NOT NULL, table_name TEXT,
        item TEXT, old_value TEXT, new_value TEXT,
        changed_by TEXT, changed_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP, note TEXT)""",
]

_MIGRATE_ANNUAL_GOALS_SQLITE = [
    """CREATE TABLE IF NOT EXISTS annual_goals (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        year INTEGER NOT NULL, dept TEXT NOT NULL, item TEXT NOT NULL,
        goal REAL DEFAULT 0,
        updated_at DATETIME DEFAULT CURRENT_TIMESTAMP,
        UNIQUE(year, dept, item))""",
]
_MIGRATE_ANNUAL_GOALS_PG = [
    """CREATE TABLE IF NOT EXISTS annual_goals (
        id SERIAL PRIMARY KEY,
        year INTEGER NOT NULL, dept TEXT NOT NULL, item TEXT NOT NULL,
        goal REAL DEFAULT 0,
        updated_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
        UNIQUE(year, dept, item))""",
]

_MIGRATE_REVENUE = [
    "ALTER TABLE revenue ADD COLUMN expected_amount REAL DEFAULT 0",
]

_MIGRATE_DEPT_GROUPS_SQLITE = [
    """CREATE TABLE IF NOT EXISTS dept_groups (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        dept TEXT NOT NULL, group_name TEXT NOT NULL, sort_order INTEGER DEFAULT 0,
        UNIQUE(dept, group_name))""",
]
_MIGRATE_DEPT_GROUPS_PG = [
    """CREATE TABLE IF NOT EXISTS dept_groups (
        id SERIAL PRIMARY KEY,
        dept TEXT NOT NULL, group_name TEXT NOT NULL, sort_order INTEGER DEFAULT 0,
        UNIQUE(dept, group_name))""",
]

_MIGRATE_CONTRACTS = [
    "ALTER TABLE contracts ADD COLUMN cross_dept INTEGER DEFAULT 0",
    "ALTER TABLE contracts ADD COLUMN cross_dept_data TEXT DEFAULT '{}'",
    "ALTER TABLE contracts ADD COLUMN payment_type TEXT DEFAULT '當年'",
    "ALTER TABLE contracts ADD COLUMN installments INTEGER DEFAULT 1",
    "ALTER TABLE contracts ADD COLUMN installment_data TEXT DEFAULT '[]'",
    "ALTER TABLE contracts ADD COLUMN group_name TEXT DEFAULT ''",
    "ALTER TABLE contracts ADD COLUMN expected_amount REAL DEFAULT 0",
    "ALTER TABLE contracts ADD COLUMN expected_date TEXT DEFAULT ''",
    "ALTER TABLE contracts ADD COLUMN project_name TEXT DEFAULT ''",
    "ALTER TABLE contracts ADD COLUMN lead_dept TEXT DEFAULT ''",
]

_OLD_DERIVE_ITEMS = ['衍生支出-研發成果', '衍生支出-研發成果(能專)', '衍生支出-其他', '衍生支出-成果下放']
_NEW_DERIVE_ITEM = '計畫衍生支出'

_OLD_SVC_INCOME = ['試驗服務收入', '技術服務收入', '訓練服務收入', '其他業務收入']
_OLD_INCOME_REMOVE = ['配合款-其他政府收入', '配合款-其他民間收入', '其他專案-來自民間']
_NEW_SVC_INCOME = '其他民間收入'

def _migrate(cur, is_pg):
    """升級舊版資料表（新增欄位）"""
    new_tables = _MIGRATE_NEW_TABLES_PG if is_pg else _MIGRATE_NEW_TABLES_SQLITE
    goal_tables = _MIGRATE_ANNUAL_GOALS_PG if is_pg else _MIGRATE_ANNUAL_GOALS_SQLITE
    dg_tables = _MIGRATE_DEPT_GROUPS_PG if is_pg else _MIGRATE_DEPT_GROUPS_SQLITE
    for stmt in _MIGRATE_USERS + _MIGRATE_REVENUE + _MIGRATE_CONTRACTS + new_tables + goal_tables + dg_tables:
        try:
            cur.execute(stmt)
        except Exception:
            pass  # 欄位/表已存在時忽略

    # 合併舊4個衍生支出項目為計畫衍生支出
    ph = '%s' if is_pg else '?'
    try:
        if is_pg:
            cur.execute("""
                INSERT INTO revenue (year, dept, month, item, amount)
                SELECT year, dept, month, %s, SUM(amount)
                FROM revenue WHERE item = ANY(%s)
                GROUP BY year, dept, month
                ON CONFLICT (year, dept, month, item) DO UPDATE SET amount = EXCLUDED.amount
            """, (_NEW_DERIVE_ITEM, _OLD_DERIVE_ITEMS))
            cur.execute("DELETE FROM revenue WHERE item = ANY(%s)", (_OLD_DERIVE_ITEMS,))
        else:
            placeholders = ','.join(['?' for _ in _OLD_DERIVE_ITEMS])
            cur.execute(f"""
                INSERT OR REPLACE INTO revenue (year, dept, month, item, amount)
                SELECT year, dept, month, ?, SUM(amount)
                FROM revenue WHERE item IN ({placeholders})
                GROUP BY year, dept, month
            """, [_NEW_DERIVE_ITEM] + _OLD_DERIVE_ITEMS)
            cur.execute(f"DELETE FROM revenue WHERE item IN ({placeholders})", _OLD_DERIVE_ITEMS)
        # 同樣合併 annual_goals
        if is_pg:
            cur.execute("""
                INSERT INTO annual_goals (year, dept, item, goal)
                SELECT year, dept, %s, SUM(goal)
                FROM annual_goals WHERE item = ANY(%s)
                GROUP BY year, dept
                ON CONFLICT (year, dept, item) DO UPDATE SET goal = EXCLUDED.goal
            """, (_NEW_DERIVE_ITEM, _OLD_DERIVE_ITEMS))
            cur.execute("DELETE FROM annual_goals WHERE item = ANY(%s)", (_OLD_DERIVE_ITEMS,))
        else:
            cur.execute(f"""
                INSERT OR REPLACE INTO annual_goals (year, dept, item, goal)
                SELECT year, dept, ?, SUM(goal)
                FROM annual_goals WHERE item IN ({placeholders})
                GROUP BY year, dept
            """, [_NEW_DERIVE_ITEM] + _OLD_DERIVE_ITEMS)
            cur.execute(f"DELETE FROM annual_goals WHERE item IN ({placeholders})", _OLD_DERIVE_ITEMS)
    except Exception:
        pass

    # 合併舊4個服務收入項目為其他民間收入
    try:
        all_old = _OLD_SVC_INCOME + _OLD_INCOME_REMOVE
        ph_list = ','.join(['%s' if is_pg else '?' for _ in _OLD_SVC_INCOME])
        all_ph = ','.join(['%s' if is_pg else '?' for _ in all_old])
        if is_pg:
            cur.execute(f"""
                INSERT INTO revenue (year, dept, month, item, amount)
                SELECT year, dept, month, %s, SUM(amount)
                FROM revenue WHERE item = ANY(%s)
                GROUP BY year, dept, month
                ON CONFLICT (year, dept, month, item) DO UPDATE SET amount = EXCLUDED.amount
            """, (_NEW_SVC_INCOME, _OLD_SVC_INCOME))
            cur.execute("DELETE FROM revenue WHERE item = ANY(%s)", (all_old,))
            cur.execute("""
                INSERT INTO annual_goals (year, dept, item, goal)
                SELECT year, dept, %s, SUM(goal)
                FROM annual_goals WHERE item = ANY(%s)
                GROUP BY year, dept
                ON CONFLICT (year, dept, item) DO UPDATE SET goal = EXCLUDED.goal
            """, (_NEW_SVC_INCOME, _OLD_SVC_INCOME))
            cur.execute("DELETE FROM annual_goals WHERE item = ANY(%s)", (all_old,))
        else:
            svc_ph = ','.join(['?' for _ in _OLD_SVC_INCOME])
            cur.execute(f"""
                INSERT OR REPLACE INTO revenue (year, dept, month, item, amount)
                SELECT year, dept, month, ?, SUM(amount)
                FROM revenue WHERE item IN ({svc_ph})
                GROUP BY year, dept, month
            """, [_NEW_SVC_INCOME] + _OLD_SVC_INCOME)
            cur.execute(f"DELETE FROM revenue WHERE item IN ({all_ph})", all_old)
            cur.execute(f"""
                INSERT OR REPLACE INTO annual_goals (year, dept, item, goal)
                SELECT year, dept, ?, SUM(goal)
                FROM annual_goals WHERE item IN ({svc_ph})
                GROUP BY year, dept
            """, [_NEW_SVC_INCOME] + _OLD_SVC_INCOME)
            cur.execute(f"DELETE FROM annual_goals WHERE item IN ({all_ph})", all_old)
    except Exception:
        pass

def init_db():
    if IS_PG:
        conn = psycopg2.connect(_DATABASE_URL)
        cur = conn.cursor()
        for stmt in _SCHEMA_PG.split(';'):
            s = stmt.strip()
            if s:
                cur.execute(s)
        _migrate(cur, True)
        cur.execute("INSERT INTO settings (key,value) VALUES ('current_year',%s) ON CONFLICT DO NOTHING",
                    (str(CURRENT_ROC_YEAR),))
        cur.execute("SELECT COUNT(*) FROM users")
        if cur.fetchone()[0] == 0:
            cur.execute("INSERT INTO users (username,password_hash,display_name,role) VALUES (%s,%s,%s,%s)",
                        ('admin', hash_pw('admin1234'), '系統管理員', 'admin'))
            print('已建立預設帳號: admin / admin1234')
        conn.commit(); conn.close()
    else:
        con = sqlite3.connect(DB)
        cur = con.cursor()
        cur.executescript(_SCHEMA_SQLITE)
        con.commit()
        _migrate(cur, False)
        cur.execute("INSERT OR IGNORE INTO settings (key,value) VALUES ('current_year',?)",
                    (str(CURRENT_ROC_YEAR),))
        cur.execute("SELECT COUNT(*) FROM users")
        if cur.fetchone()[0] == 0:
            cur.execute("INSERT INTO users (username,password_hash,display_name,role) VALUES (?,?,?,?)",
                        ('admin', hash_pw('admin1234'), '系統管理員', 'admin'))
            print('已建立預設帳號: admin / admin1234')
        con.commit(); con.close()

def get_db():
    if IS_PG:
        return _PgWrapper(psycopg2.connect(_DATABASE_URL))
    con = sqlite3.connect(DB)
    con.row_factory = sqlite3.Row
    return con

def get_current_year():
    """取得目前選用的年度（從 session 或 DB settings）"""
    if session.get('year'):
        return int(session['year'])
    con = get_db()
    row = con.execute("SELECT value FROM settings WHERE key='current_year'").fetchone()
    con.close()
    return int(row['value']) if row else CURRENT_ROC_YEAR

def get_user_role():
    return session.get('role', 'viewer')

def is_admin():
    return get_user_role() == 'admin'

def can_write():
    """admin 或 editor 才能寫入。"""
    return get_user_role() in ('admin', 'editor')

def get_allowed_depts():
    """依使用者角色/部門傳回可存取（可看）的部門清單。"""
    if is_admin():
        return DEPARTMENTS
    dept = session.get('dept', '')
    if dept and dept in DEPARTMENTS:
        return [dept]
    return DEPARTMENTS  # 部門不在清單 → 全部可看

def can_access_dept(dept):
    """是否能讀取指定部門資料。"""
    if is_admin():
        return True
    user_dept = session.get('dept', '')
    if not user_dept or user_dept not in DEPARTMENTS:
        return True
    return user_dept == dept

def can_write_dept(dept):
    """是否能寫入指定部門資料（admin 全部；editor 限自己部門）。"""
    if not can_write():
        return False
    if is_admin():
        return True
    user_dept = session.get('dept', '')
    if not user_dept or user_dept not in DEPARTMENTS:
        return True
    return user_dept == dept

def is_admin_or_no_dept():
    if is_admin():
        return True
    dept = session.get('dept', '')
    return not dept or dept not in DEPARTMENTS

def require_write(dept=None):
    """若無寫入權限則回傳 403，否則 None。dept 不為 None 時同時檢查部門。"""
    if dept is not None:
        if not can_write_dept(dept):
            return jsonify({'error': '無權限修改此部門資料'}), 403
    elif not can_write():
        return jsonify({'error': '無寫入權限'}), 403
    return None

def get_all_years():
    """取得資料庫中有資料的所有年度（+ 當前年度）"""
    con = get_db()
    years = set()
    for tbl in ('revenue', 'contracts', 'unclaimed'):
        rows = con.execute(f"SELECT DISTINCT year FROM {tbl}").fetchall()
        years.update(r['year'] for r in rows)
    cur_yr = get_current_year()
    years.add(cur_yr)
    # 加入前後各 1 年供切換
    years.add(cur_yr - 1)
    years.add(cur_yr + 1)
    con.close()
    return sorted(years, reverse=True)

# ── 帳號路由 ───────────────────────────────────────────
@app.route('/login', methods=['GET', 'POST'])
def login():
    if session.get('user'):
        return redirect(url_for('index'))
    error = None
    if request.method == 'POST':
        username = request.form.get('username', '').strip()
        password = request.form.get('password', '')
        con = get_db()
        user = con.execute(
            "SELECT * FROM users WHERE username=? AND password_hash=?",
            (username, hash_pw(password))
        ).fetchone()
        con.close()
        if user:
            if user['disabled']:
                error = '此帳號已被停用，請聯繫管理員。'
                return render_template('login.html', error=error)
            session.permanent = False
            session['user'] = user['username']
            session['display_name'] = user['display_name'] or user['username']
            session['role'] = user['role']
            session['dept'] = user['dept'] or ''
            # 預設使用 DB 設定年度
            con2 = get_db()
            row = con2.execute("SELECT value FROM settings WHERE key='current_year'").fetchone()
            con2.close()
            session['year'] = int(row['value']) if row else CURRENT_ROC_YEAR
            return redirect(request.args.get('next') or url_for('index'))
        error = '帳號或密碼錯誤，請再試一次。'
    return render_template('login.html', error=error)

@app.route('/logout')
def logout():
    session.clear()
    return redirect(url_for('login'))

@app.route('/api/switch_year', methods=['POST'])
@login_required
def switch_year():
    year = request.json.get('year')
    if year:
        session['year'] = int(year)
    return jsonify({'status': 'ok', 'year': session.get('year')})

@app.route('/forgot_password', methods=['GET', 'POST'])
def forgot_password():
    msg = None
    token_shown = None
    if request.method == 'POST':
        username = request.form.get('username', '').strip()
        con = get_db()
        user = con.execute("SELECT * FROM users WHERE username=?", (username,)).fetchone()
        if user:
            token = secrets.token_urlsafe(24)
            expires = (datetime.now() + timedelta(hours=1)).strftime('%Y-%m-%d %H:%M:%S')
            con.execute("UPDATE users SET reset_token=?, reset_expires=? WHERE username=?",
                        (token, expires, username))
            con.commit()
            token_shown = token
            msg = '重設連結已產生（有效 1 小時）。'
        else:
            msg = '找不到此帳號，請確認帳號名稱。'
        con.close()
    return render_template('forgot_password.html', msg=msg, token_shown=token_shown)

@app.route('/reset_password/<token>', methods=['GET', 'POST'])
def reset_password(token):
    con = get_db()
    now = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
    user = con.execute(
        "SELECT * FROM users WHERE reset_token=? AND reset_expires > ?", (token, now)
    ).fetchone()
    if not user:
        con.close()
        return render_template('reset_password.html', error='連結無效或已過期。')
    error = None
    if request.method == 'POST':
        pw1 = request.form.get('password', '')
        pw2 = request.form.get('password2', '')
        if len(pw1) < 6:
            error = '密碼至少需要 6 個字元。'
        elif pw1 != pw2:
            error = '兩次輸入的密碼不一致。'
        else:
            con.execute(
                "UPDATE users SET password_hash=?, reset_token=NULL, reset_expires=NULL WHERE id=?",
                (hash_pw(pw1), user['id'])
            )
            con.commit()
            con.close()
            return render_template('reset_password.html', success=True)
    con.close()
    return render_template('reset_password.html', token=token,
                           username=user['username'], error=error)

@app.route('/change_password', methods=['GET', 'POST'])
@login_required
def change_password():
    error = None
    success = False
    if request.method == 'POST':
        old_pw  = request.form.get('old_password', '')
        new_pw  = request.form.get('new_password', '')
        new_pw2 = request.form.get('new_password2', '')
        con = get_db()
        user = con.execute(
            "SELECT * FROM users WHERE username=? AND password_hash=?",
            (session['user'], hash_pw(old_pw))
        ).fetchone()
        if not user:
            error = '目前密碼不正確。'
        elif len(new_pw) < 6:
            error = '新密碼至少需要 6 個字元。'
        elif new_pw != new_pw2:
            error = '兩次輸入的新密碼不一致。'
        else:
            con.execute("UPDATE users SET password_hash=? WHERE username=?",
                        (hash_pw(new_pw), session['user']))
            con.commit()
            success = True
        con.close()
    year = get_current_year()
    return render_template('change_password.html', error=error, success=success, year=year)

@app.route('/admin/unlock_requests_page')
@login_required
def admin_unlock_requests_page():
    if session.get('role') != 'admin':
        return redirect(url_for('index'))
    year = get_current_year()
    all_years = get_all_years()
    con = get_db()
    reqs = [dict(r) for r in con.execute(
        'SELECT * FROM locks WHERE year=? AND unlock_requested=1 ORDER BY req_at DESC', (year,)
    ).fetchall()]
    con.close()
    return render_template('admin_unlock_requests.html', requests=reqs,
                           year=year, all_years=all_years)

@app.route('/admin/users')
@login_required
def admin_users():
    if session.get('role') != 'admin':
        return redirect(url_for('index'))
    con = get_db()
    users = [dict(r) for r in con.execute(
        "SELECT id, username, display_name, dept, role, disabled, created_at FROM users ORDER BY id"
    ).fetchall()]
    con.close()
    year = get_current_year()
    return render_template('admin_users.html', users=users, year=year)

@app.route('/admin/add_user', methods=['POST'])
@login_required
def add_user():
    if session.get('role') != 'admin':
        return jsonify({'error': 'forbidden'}), 403
    d = request.json
    username = d.get('username', '').strip()
    password = d.get('password', '')
    if not username or len(password) < 6:
        return jsonify({'error': '帳號或密碼格式不正確'}), 400
    con = get_db()
    try:
        con.execute("INSERT INTO users (username, password_hash, display_name, role, dept) VALUES (?,?,?,?,?)",
                    (username, hash_pw(password), d.get('display_name',''), d.get('role','user'), d.get('dept','')))
        con.commit()
    except Exception as e:
        try: con.close()
        except Exception: pass
        if 'unique' in str(e).lower() or 'integrity' in str(e).lower() or 'duplicate' in str(e).lower():
            return jsonify({'error': '帳號已存在'}), 400
        return jsonify({'error': f'新增失敗：{e}'}), 500
    con.close()
    return jsonify({'status': 'ok'})

@app.route('/admin/toggle_user/<int:uid>', methods=['POST'])
@login_required
def toggle_user(uid):
    if session.get('role') != 'admin':
        return jsonify({'error': 'forbidden'}), 403
    con = get_db()
    user = con.execute("SELECT username, disabled FROM users WHERE id=?", (uid,)).fetchone()
    if not user:
        con.close()
        return jsonify({'error': '帳號不存在'}), 404
    if user['username'] == 'admin':
        con.close()
        return jsonify({'error': '不可停用 admin 帳號'}), 400
    new_state = 0 if user['disabled'] else 1
    con.execute("UPDATE users SET disabled=? WHERE id=?", (new_state, uid))
    con.commit()
    con.close()
    return jsonify({'status': 'ok', 'disabled': new_state})

@app.route('/admin/export_users')
@login_required
def export_users():
    if session.get('role') != 'admin':
        return redirect(url_for('index'))
    con = get_db()
    users = [dict(r) for r in con.execute(
        "SELECT username, display_name, dept, role, disabled, created_at FROM users ORDER BY id"
    ).fetchall()]
    con.close()

    wb = openpyxl.Workbook()
    thin   = Side(style='thin')
    border = Border(left=thin, right=thin, top=thin, bottom=thin)
    hfill  = PatternFill('solid', start_color='BDD7EE', fgColor='BDD7EE')
    hfont  = Font(name='微軟正黑體', size=11, bold=True)
    ctr    = Alignment(horizontal='center', vertical='center')

    ws = wb.active
    ws.title = '帳號清單'
    headers = ['帳號', '姓名', '部門', '權限', '狀態', '建立時間']
    widths  = [18, 16, 14, 14, 10, 22]
    for col, (h, w) in enumerate(zip(headers, widths), 1):
        c = ws.cell(row=1, column=col, value=h)
        c.font = hfont; c.fill = hfill; c.alignment = ctr; c.border = border
        ws.column_dimensions[get_column_letter(col)].width = w

    role_map = {'admin': '管理員', 'user': '一般使用者'}
    for ri, u in enumerate(users, 2):
        vals = [
            u['username'],
            u['display_name'] or '',
            u['dept'] or '',
            role_map.get(u['role'], u['role']),
            '停用' if u['disabled'] else '啟用',
            str(u['created_at'] or '')[:10],
        ]
        for ci, v in enumerate(vals, 1):
            cell = ws.cell(row=ri, column=ci, value=v)
            cell.border = border
            if u['disabled']:
                cell.font = Font(name='微軟正黑體', size=10, color='999999')

    output = io.BytesIO()
    wb.save(output); output.seek(0)
    return send_file(output, as_attachment=True,
                     download_name=f'帳號清單_{datetime.now().strftime("%Y%m%d")}.xlsx',
                     mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')

@app.route('/admin/reset_user_password/<int:uid>', methods=['POST'])
@login_required
def reset_user_password(uid):
    if session.get('role') != 'admin':
        return jsonify({'error': 'forbidden'}), 403
    pw = request.json.get('password', '')
    if len(pw) < 6:
        return jsonify({'error': '密碼至少需要 6 個字元'}), 400
    con = get_db()
    con.execute("UPDATE users SET password_hash=? WHERE id=?", (hash_pw(pw), uid))
    con.commit()
    con.close()
    return jsonify({'status': 'ok'})

@app.route('/admin/update_user/<int:uid>', methods=['POST'])
@login_required
def update_user(uid):
    if session.get('role') != 'admin':
        return jsonify({'error': 'forbidden'}), 403
    d = request.json
    display_name = d.get('display_name', '').strip()
    dept = d.get('dept', '').strip()
    role = d.get('role', 'editor')
    con = get_db()
    con.execute("UPDATE users SET display_name=?, dept=?, role=? WHERE id=?",
                (display_name, dept, role, uid))
    con.commit()
    con.close()
    return jsonify({'status': 'ok', 'display_name': display_name, 'dept': dept, 'role': role})

@app.route('/admin/delete_user/<int:uid>', methods=['DELETE'])
@login_required
def delete_user(uid):
    if session.get('role') != 'admin':
        return jsonify({'error': 'forbidden'}), 403
    con = get_db()
    user = con.execute("SELECT username FROM users WHERE id=?", (uid,)).fetchone()
    if user and user['username'] == 'admin':
        con.close()
        return jsonify({'error': '不可刪除 admin 帳號'}), 400
    con.execute("DELETE FROM users WHERE id=?", (uid,))
    con.commit()
    con.close()
    return jsonify({'status': 'ok'})

# ── 部門組別管理 ────────────────────────────────────────
@app.route('/admin/dept_groups')
@login_required
def admin_dept_groups():
    if session.get('role') != 'admin':
        return redirect(url_for('index'))
    con = get_db()
    rows = [dict(r) for r in con.execute(
        "SELECT * FROM dept_groups ORDER BY dept, sort_order, group_name"
    ).fetchall()]
    con.close()
    # 整理成 {dept: [groups]}
    by_dept = {d: [] for d in DEPARTMENTS}
    for r in rows:
        if r['dept'] in by_dept:
            by_dept[r['dept']].append(r)
    return render_template('admin_dept_groups.html',
                           departments=DEPARTMENTS, by_dept=by_dept, year=get_current_year())

@app.route('/api/dept_groups', methods=['GET'])
@login_required
def api_get_dept_groups():
    dept = request.args.get('dept', '')
    con = get_db()
    if dept:
        rows = [dict(r) for r in con.execute(
            "SELECT id, group_name, sort_order FROM dept_groups WHERE dept=? ORDER BY sort_order, group_name",
            (dept,)
        ).fetchall()]
    else:
        rows = [dict(r) for r in con.execute(
            "SELECT * FROM dept_groups ORDER BY dept, sort_order, group_name"
        ).fetchall()]
    con.close()
    return jsonify(rows)

@app.route('/api/dept_groups', methods=['POST'])
@login_required
def api_add_dept_group():
    if session.get('role') != 'admin':
        return jsonify({'error': 'forbidden'}), 403
    d = request.json
    dept = d.get('dept', '').strip()
    group_name = d.get('group_name', '').strip()
    sort_order = d.get('sort_order', 0)
    if not dept or not group_name:
        return jsonify({'error': '部門和組別名稱必填'}), 400
    if dept not in DEPARTMENTS:
        return jsonify({'error': '無效的部門'}), 400
    con = get_db()
    ph = '%s' if IS_PG else '?'
    try:
        con.execute(f"INSERT INTO dept_groups (dept, group_name, sort_order) VALUES ({ph},{ph},{ph})",
                    (dept, group_name, sort_order))
        con.commit()
        row = con.execute(f"SELECT * FROM dept_groups WHERE dept={ph} AND group_name={ph}",
                          (dept, group_name)).fetchone()
        con.close()
        return jsonify(dict(row))
    except Exception as e:
        con.close()
        return jsonify({'error': '組別已存在或發生錯誤: ' + str(e)}), 400

@app.route('/api/dept_groups/<int:gid>', methods=['DELETE'])
@login_required
def api_delete_dept_group(gid):
    if session.get('role') != 'admin':
        return jsonify({'error': 'forbidden'}), 403
    ph = '%s' if IS_PG else '?'
    con = get_db()
    con.execute(f"DELETE FROM dept_groups WHERE id={ph}", (gid,))
    con.commit()
    con.close()
    return jsonify({'status': 'ok'})

@app.route('/api/dept_groups/<int:gid>', methods=['PUT'])
@login_required
def api_update_dept_group(gid):
    if session.get('role') != 'admin':
        return jsonify({'error': 'forbidden'}), 403
    d = request.json
    group_name = d.get('group_name', '').strip()
    sort_order = d.get('sort_order', 0)
    if not group_name:
        return jsonify({'error': '組別名稱必填'}), 400
    ph = '%s' if IS_PG else '?'
    con = get_db()
    con.execute(f"UPDATE dept_groups SET group_name={ph}, sort_order={ph} WHERE id={ph}",
                (group_name, sort_order, gid))
    con.commit()
    con.close()
    return jsonify({'status': 'ok'})

# ── 主功能路由 ─────────────────────────────────────────
@app.route('/')
@login_required
def index():
    year = get_current_year()
    all_years = get_all_years()
    allowed = get_allowed_depts()
    # 有部門限制的使用者直接跳到自己部門
    if len(allowed) == 1:
        return redirect(url_for('dept_view', dept=allowed[0],
                                month=datetime.now().month))
    return render_template('index.html', departments=allowed, months=MONTHS,
                           year=year, all_years=all_years)

@app.route('/dept/<dept>')
@login_required
def dept_view(dept):
    if dept not in DEPARTMENTS:
        return redirect(url_for('index'))
    if not can_access_dept(dept):
        return redirect(url_for('dept_view', dept=session.get('dept')))
    month = request.args.get('month', 1, type=int)
    year = get_current_year()
    all_years = get_all_years()
    allowed = get_allowed_depts()
    return render_template('dept.html', dept=dept, month=month,
                           months=MONTHS, year=year, all_years=all_years,
                           departments=allowed,
                           income_items=INCOME_ITEMS,
                           expense_items=EXPENSE_ITEMS,
                           unclaimed_items=UNCLAIMED_ITEMS,
                           can_write=can_write_dept(dept))

@app.route('/api/data/<dept>/<int:month>')
@login_required
def get_data(dept, month):
    if not can_access_dept(dept):
        return jsonify({'error': 'forbidden'}), 403
    year = request.args.get('year', type=int) or get_current_year()
    con = get_db()
    rows = con.execute(
        'SELECT item, amount, expected_amount, goal FROM revenue WHERE year=? AND dept=? AND month=?',
        (year, dept, month)
    ).fetchall()
    data = {r['item']: {'amount': r['amount'], 'expected_amount': r['expected_amount'] or 0, 'goal': r['goal']} for r in rows}
    cumul_data = {r['item']: r['amount'] for r in rows}
    unclaimed = con.execute(
        'SELECT item, amount FROM unclaimed WHERE year=? AND dept=? AND month=?',
        (year, dept, month)
    ).fetchall()
    unclaim_data = {r['item']: r['amount'] for r in unclaimed}
    contracts = [dict(r) for r in con.execute(
        'SELECT * FROM contracts WHERE year=? AND dept=? AND month=?',
        (year, dept, month)
    ).fetchall()]
    carry_forward = [dict(r) for r in con.execute(
        'SELECT * FROM contracts WHERE year=? AND dept=? AND month=? AND carry_next=1',
        (year, dept, month - 1)
    ).fetchall()] if month > 1 else []
    con.close()
    return jsonify({'revenue': data, 'cumul': cumul_data, 'unclaimed': unclaim_data,
                    'contracts': contracts, 'carry_forward': carry_forward})

def _is_locked(con, year, dept, month):
    row = con.execute('SELECT locked FROM locks WHERE year=? AND dept=? AND month=?',
                      (year, dept, month)).fetchone()
    return bool(row and row['locked'])

def _log(con, year, dept, month, action, table_name='', item='', old_val='', new_val='', note=''):
    con.execute('''INSERT INTO audit_log (year,dept,month,action,table_name,item,old_value,new_value,changed_by,note)
        VALUES (?,?,?,?,?,?,?,?,?,?)''',
        (year, dept, month, action, table_name, item, str(old_val), str(new_val),
         session.get('user','?'), note))

@app.route('/api/lock_status/<dept>/<int:month>')
@login_required
def lock_status(dept, month):
    year = get_current_year()
    con = get_db()
    row = con.execute('SELECT * FROM locks WHERE year=? AND dept=? AND month=?',
                      (year, dept, month)).fetchone()
    con.close()
    if not row:
        return jsonify({'locked': False, 'unlock_requested': False})
    return jsonify({
        'locked': bool(row['locked']),
        'locked_by': row['locked_by'] or '',
        'locked_at': str(row['locked_at'] or ''),
        'unlock_requested': bool(row['unlock_requested']),
        'req_by': row['req_by'] or '',
        'req_reason': row['req_reason'] or '',
    })

@app.route('/api/lock_dept', methods=['POST'])
@login_required
def lock_dept():
    d = request.json
    year = get_current_year()
    dept, month = d['dept'], d['month']
    if not can_write_dept(dept):
        return jsonify({'error': '無權限鎖定此部門'}), 403
    now = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
    con = get_db()
    if _is_locked(con, year, dept, month):
        con.close()
        return jsonify({'error': '已鎖定'}), 400
    con.execute('''INSERT INTO locks (year,dept,month,locked,locked_by,locked_at)
        VALUES (?,?,?,1,?,?)
        ON CONFLICT(year,dept,month) DO UPDATE
        SET locked=1, locked_by=excluded.locked_by, locked_at=excluded.locked_at,
            unlock_requested=0, req_by=NULL, req_at=NULL, unlocked_by=NULL, unlocked_at=NULL''',
        (year, dept, month, session['user'], now))
    _log(con, year, dept, month, 'lock', note=f'由 {session["user"]} 鎖定')
    con.commit(); con.close()
    return jsonify({'status': 'ok'})

@app.route('/api/request_unlock', methods=['POST'])
@login_required
def request_unlock():
    d = request.json
    year = get_current_year()
    dept, month = d['dept'], d['month']
    if not can_write_dept(dept):
        return jsonify({'error': '無權限申請解鎖'}), 403
    reason = d.get('reason', '').strip()
    now = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
    con = get_db()
    con.execute('''INSERT INTO locks (year,dept,month,locked,unlock_requested,req_by,req_at,req_reason)
        VALUES (?,?,?,1,1,?,?,?)
        ON CONFLICT(year,dept,month) DO UPDATE
        SET unlock_requested=1, req_by=excluded.req_by, req_at=excluded.req_at,
            req_reason=excluded.req_reason''',
        (year, dept, month, session['user'], now, reason))
    _log(con, year, dept, month, 'request_unlock', note=f'申請解鎖，原因：{reason}')
    con.commit(); con.close()
    return jsonify({'status': 'ok'})

@app.route('/api/unlock_dept', methods=['POST'])
@login_required
def unlock_dept():
    if session.get('role') != 'admin':
        return jsonify({'error': 'forbidden'}), 403
    d = request.json
    year = get_current_year()
    dept, month = d['dept'], d['month']
    now = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
    con = get_db()
    con.execute('''UPDATE locks SET locked=0, unlock_requested=0,
        unlocked_by=?, unlocked_at=? WHERE year=? AND dept=? AND month=?''',
        (session['user'], now, year, dept, month))
    _log(con, year, dept, month, 'unlock', note=f'由管理員 {session["user"]} 解鎖')
    con.commit(); con.close()
    return jsonify({'status': 'ok'})

@app.route('/api/reject_unlock', methods=['POST'])
@login_required
def reject_unlock():
    if session.get('role') != 'admin':
        return jsonify({'error': 'forbidden'}), 403
    d = request.json
    year = get_current_year()
    dept, month = d['dept'], d['month']
    con = get_db()
    con.execute('''UPDATE locks SET unlock_requested=0, req_by=NULL, req_at=NULL, req_reason=NULL
        WHERE year=? AND dept=? AND month=?''', (year, dept, month))
    _log(con, year, dept, month, 'reject_unlock', note=f'管理員 {session["user"]} 拒絕解鎖申請')
    con.commit(); con.close()
    return jsonify({'status': 'ok'})

@app.route('/api/unlock_requests')
@login_required
def unlock_requests():
    if session.get('role') != 'admin':
        return jsonify({'error': 'forbidden'}), 403
    year = get_current_year()
    con = get_db()
    rows = [dict(r) for r in con.execute(
        'SELECT * FROM locks WHERE year=? AND unlock_requested=1 ORDER BY req_at DESC', (year,)
    ).fetchall()]
    con.close()
    return jsonify({'requests': rows})

@app.route('/api/audit_log/<dept>/<int:month>')
@login_required
def get_audit_log(dept, month):
    if not can_access_dept(dept) and session.get('role') != 'admin':
        return jsonify({'error': 'forbidden'}), 403
    year = get_current_year()
    con = get_db()
    rows = [dict(r) for r in con.execute(
        '''SELECT * FROM audit_log WHERE year=? AND dept=? AND month=?
           ORDER BY changed_at DESC LIMIT 100''',
        (year, dept, month)
    ).fetchall()]
    con.close()
    return jsonify({'logs': rows})

@app.route('/api/save_revenue', methods=['POST'])
@login_required
def save_revenue():
    d = request.json
    year = get_current_year()
    dept, month = d['dept'], d['month']
    if not can_write_dept(dept):
        return jsonify({'error': '無寫入權限'}), 403
    con = get_db()
    if _is_locked(con, year, dept, month):
        con.close()
        return jsonify({'error': 'locked'}), 423
    for item, vals in d['items'].items():
        old = con.execute('SELECT amount,expected_amount,goal FROM revenue WHERE year=? AND dept=? AND month=? AND item=?',
                          (year,dept,month,item)).fetchone()
        new_amt = vals.get('amount', 0)
        new_exp = vals.get('expected_amount', 0)
        new_goal = vals.get('goal', 0)
        if old and (old['amount'] != new_amt or old['goal'] != new_goal):
            _log(con, year, dept, month, 'edit', 'revenue', item,
                 f"金額:{old['amount']},目標:{old['goal']}",
                 f"金額:{new_amt},目標:{new_goal}")
        con.execute('''INSERT INTO revenue (year, dept, month, item, amount, expected_amount, goal)
            VALUES (?,?,?,?,?,?,?)
            ON CONFLICT(year, dept, month, item) DO UPDATE
            SET amount=excluded.amount, expected_amount=excluded.expected_amount,
                goal=excluded.goal, updated_at=CURRENT_TIMESTAMP''',
            (year, dept, month, item, new_amt, new_exp, new_goal))
    con.commit()
    con.close()
    return jsonify({'status': 'ok'})

@app.route('/api/annual_goals/<dept>')
@login_required
def get_annual_goals(dept):
    if not can_access_dept(dept):
        return jsonify({'error': 'forbidden'}), 403
    year = request.args.get('year', type=int) or get_current_year()
    con = get_db()
    rows = con.execute('SELECT item, goal FROM annual_goals WHERE year=? AND dept=?', (year, dept)).fetchall()
    con.close()
    return jsonify({r['item']: r['goal'] for r in rows})

@app.route('/api/import_goals_excel', methods=['POST'])
@login_required
def import_goals_excel():
    dept = request.form.get('dept', '')
    if not can_write_dept(dept):
        return jsonify({'error': '無寫入權限'}), 403
    f = request.files.get('file')
    if not f:
        return jsonify({'error': 'no file'}), 400
    try:
        wb = openpyxl.load_workbook(io.BytesIO(f.read()), data_only=True)
    except Exception as e:
        return jsonify({'error': f'無法讀取 Excel：{e}'}), 400

    all_items = INCOME_ITEMS + EXPENSE_ITEMS
    # strip known suffixes for matching
    def _norm(s):
        return str(s or '').strip().rstrip('◎').strip()

    item_set = {_norm(i): i for i in all_items}
    # Excel 舊名稱對照
    _aliases = {'計畫衍生收入': '科專衍生收入'}
    for old, new in _aliases.items():
        if new in [v for v in item_set.values()]:
            item_set[old] = new
    goals = {}

    for sh_name in wb.sheetnames:
        ws = wb[sh_name]
        # find which dept this sheet belongs to
        sheet_dept = None
        goal_col = None   # 0-based column index of 「今年年度預算數」
        for row in ws.iter_rows(min_row=1, max_row=10, values_only=True):
            for cell in row:
                if cell and '部門：' in str(cell):
                    sheet_dept = str(cell).replace('部門：', '').strip()
                    break
            # find header row to determine column
            row_vals = [str(v or '') for v in row]
            count = sum(1 for v in row_vals if '年度預算數' in v)
            if count >= 2:
                # second occurrence is the current-year goal column
                found = 0
                for ci, v in enumerate(row_vals):
                    if '年度預算數' in v:
                        found += 1
                        if found == 2:
                            goal_col = ci
                            break

        if sheet_dept is None or goal_col is None:
            continue
        if dept not in sheet_dept and sheet_dept not in dept:
            continue

        # extract item → goal
        for row in ws.iter_rows(values_only=True):
            name = _norm(row[0] if row else None)
            if name in item_set:
                try:
                    val = float(row[goal_col] or 0)
                except (TypeError, ValueError):
                    val = 0
                goals[item_set[name]] = val

    if not goals:
        return jsonify({'error': f'未在 Excel 中找到「{dept}」部門資料'}), 404
    return jsonify({'goals': goals})

@app.route('/api/save_annual_goals', methods=['POST'])
@login_required
def save_annual_goals():
    d = request.json
    dept = d.get('dept', '')
    if not can_write_dept(dept):
        return jsonify({'error': '無寫入權限'}), 403
    year = int(d.get('year') or get_current_year())
    con = get_db()
    for item, goal in d.get('goals', {}).items():
        g = float(goal or 0)
        con.execute('''INSERT INTO annual_goals (year, dept, item, goal)
            VALUES (?,?,?,?)
            ON CONFLICT(year, dept, item) DO UPDATE
            SET goal=excluded.goal, updated_at=CURRENT_TIMESTAMP''',
            (year, dept, item, g))
    con.commit()
    con.close()
    return jsonify({'status': 'ok'})

@app.route('/api/save_unclaimed', methods=['POST'])
@login_required
def save_unclaimed():
    d = request.json
    year = get_current_year()
    dept, month = d['dept'], d['month']
    if not can_write_dept(dept):
        return jsonify({'error': '無寫入權限'}), 403
    con = get_db()
    if _is_locked(con, year, dept, month):
        con.close()
        return jsonify({'error': 'locked'}), 423
    for item, amount in d['items'].items():
        old = con.execute('SELECT amount FROM unclaimed WHERE year=? AND dept=? AND month=? AND item=?',
                          (year,dept,month,item)).fetchone()
        if old and old['amount'] != amount:
            _log(con, year, dept, month, 'edit', 'unclaimed', item, old['amount'], amount)
        con.execute('''INSERT INTO unclaimed (year, dept, month, item, amount)
            VALUES (?,?,?,?,?)
            ON CONFLICT(year, dept, month, item) DO UPDATE
            SET amount=excluded.amount, updated_at=CURRENT_TIMESTAMP''',
            (year, dept, month, item, amount))
    con.commit()
    con.close()
    return jsonify({'status': 'ok'})

@app.route('/api/save_contract', methods=['POST'])
@login_required
def save_contract():
    import json as _json
    d = request.json
    year = get_current_year()
    if not can_write_dept(d.get('dept', '')):
        return jsonify({'error': '無寫入權限'}), 403
    con = get_db()
    if _is_locked(con, year, d.get('dept',''), d.get('month',0)):
        con.close()
        return jsonify({'error': 'locked'}), 423
    cross_dept_data = _json.dumps(d.get('cross_dept_data', {}), ensure_ascii=False)
    installment_data = _json.dumps(d.get('installment_data', []), ensure_ascii=False)
    expected_amount = d.get('expected_amount', 0)
    expected_date = d.get('expected_date', '')
    project_name = d.get('project_name', '')
    lead_dept = d.get('lead_dept', '')
    if d.get('id'):
        con.execute('''UPDATE contracts SET client=?, project_name=?, amount=?, sign_date=?,
            status=?, group_name=?, note=?, carry_next=?,
            cross_dept=?, cross_dept_data=?,
            payment_type=?, installments=?, installment_data=?,
            expected_amount=?, expected_date=?, lead_dept=?,
            updated_at=CURRENT_TIMESTAMP
            WHERE id=?''',
            (d['client'], project_name, d['amount'], d.get('sign_date',''),
             d['status'], d.get('group_name',''), d.get('note',''), d.get('carry_next',0),
             1 if d.get('cross_dept') else 0, cross_dept_data,
             d.get('payment_type','當年'), d.get('installments',1), installment_data,
             expected_amount, expected_date, lead_dept,
             d['id']))
    else:
        con.execute('''INSERT INTO contracts
            (year, dept, month, client, project_name, amount, sign_date,
             status, group_name, note, carry_next,
             cross_dept, cross_dept_data,
             payment_type, installments, installment_data,
             expected_amount, expected_date, lead_dept)
            VALUES (?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?)''',
            (year, d['dept'], d['month'], d['client'], project_name, d['amount'],
             d.get('sign_date',''), d['status'],
             d.get('group_name',''), d.get('note',''), d.get('carry_next',0),
             1 if d.get('cross_dept') else 0, cross_dept_data,
             d.get('payment_type','當年'), d.get('installments',1), installment_data,
             expected_amount, expected_date, lead_dept))
    con.commit()
    if not d.get('id'):
        new_id = con.execute('SELECT last_insert_rowid()').fetchone()[0] if not IS_PG else \
                 con.execute('SELECT lastval()').fetchone()[0]
        con.close()
        return jsonify({'status': 'ok', 'id': new_id})
    con.close()
    return jsonify({'status': 'ok'})

@app.route('/api/export_dept_excel/<dept>/<int:month>')
@login_required
def export_dept_excel(dept, month):
    if not can_access_dept(dept):
        return jsonify({'error': 'forbidden'}), 403
    year = get_current_year()
    con = get_db()
    rev_rows     = con.execute('SELECT item,amount,goal FROM revenue WHERE year=? AND dept=? AND month=? ORDER BY item', (year,dept,month)).fetchall()
    unclaim_rows = con.execute('SELECT item,amount FROM unclaimed WHERE year=? AND dept=? AND month=?', (year,dept,month)).fetchall()
    con.close()

    wb = openpyxl.Workbook()
    thin = Side(style='thin'); border = Border(left=thin,right=thin,top=thin,bottom=thin)
    hfill = PatternFill('solid',start_color='BDD7EE',fgColor='BDD7EE')
    hfont = Font(name='微軟正黑體',size=10,bold=True)
    ctr   = Alignment(horizontal='center',vertical='center')

    def hdr(ws, cols, widths):
        for c,(h,w) in enumerate(zip(cols,widths),1):
            cell = ws.cell(1,c,h); cell.font=hfont; cell.fill=hfill
            cell.alignment=ctr; cell.border=border
            ws.column_dimensions[get_column_letter(c)].width=w

    ws1 = wb.active; ws1.title = f'{month}月收支'
    ws1['A1'] = f'{year}年 {dept} {month}月 來自民間業務收支表'
    ws1['A1'].font = Font(name='微軟正黑體',size=12,bold=True)
    ws1.merge_cells('A1:D1'); ws1.row_dimensions[1].height = 22
    hdr_row = ['項目','金額(元)','年度目標','達成率']
    for c,(h,w) in enumerate(zip(hdr_row,[28,14,14,10]),1):
        cell = ws1.cell(2,c,h); cell.font=hfont; cell.fill=hfill
        cell.alignment=ctr; cell.border=border
        ws1.column_dimensions[get_column_letter(c)].width=w
    uc_map = {r['item']:r['amount'] for r in unclaim_rows}
    for ri,r in enumerate(rev_rows,3):
        g = r['goal'] or 0; a = r['amount'] or 0
        ws1.cell(ri,1,r['item']).border=border
        ws1.cell(ri,2,a).border=border
        ws1.cell(ri,3,g).border=border
        ws1.cell(ri,4, f'{a/g*100:.1f}%' if g else '-').border=border
    if uc_map:
        ri = len(rev_rows)+3
        ws1.cell(ri,1,'── 已申請未核銷 ──').font=Font(name='微軟正黑體',bold=True,color='C00000')
        for item,amt in uc_map.items():
            ri+=1
            ws1.cell(ri,1,item).border=border
            ws1.cell(ri,2,amt).border=border

    output = io.BytesIO(); wb.save(output); output.seek(0)
    return send_file(output, as_attachment=True,
                     download_name=f'{year}年{dept}{month}月收支.xlsx',
                     mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')

@app.route('/api/export_contracts_excel/<dept>/<int:month>')
@login_required
def export_contracts_excel(dept, month):
    if not can_access_dept(dept):
        return jsonify({'error': 'forbidden'}), 403
    import json as _json
    year = get_current_year()
    con = get_db()
    rows = [dict(r) for r in con.execute(
        'SELECT * FROM contracts WHERE year=? AND dept=? AND month=? ORDER BY id',
        (year,dept,month)).fetchall()]
    con.close()

    wb = openpyxl.Workbook()
    thin = Side(style='thin'); border = Border(left=thin,right=thin,top=thin,bottom=thin)
    hfill = PatternFill('solid',start_color='BDD7EE',fgColor='BDD7EE')
    hfont = Font(name='微軟正黑體',size=10,bold=True)
    ctr   = Alignment(horizontal='center',vertical='center')

    ws = wb.active; ws.title = f'{month}月合約'
    ws['A1'] = f'{year}年 {dept} {month}月 合約追蹤'
    ws['A1'].font = Font(name='微軟正黑體',size=12,bold=True)
    ws.merge_cells('A1:M1'); ws.row_dimensions[1].height = 22
    headers = ['洽談廠商/客戶','計畫名稱','組別','狀態','預計簽約金額','預計簽約日期','簽約金額','簽約日期','金額方式','期數','跨部門','延續下月','備註']
    widths  = [22,28,12,14,16,14,14,12,10,8,20,8,20]
    for c,(h,w) in enumerate(zip(headers,widths),1):
        cell = ws.cell(2,c,h); cell.font=hfont; cell.fill=hfill
        cell.alignment=ctr; cell.border=border
        ws.column_dimensions[get_column_letter(c)].width=w
    for ri,r in enumerate(rows,3):
        cd = _json.loads(r.get('cross_dept_data') or '{}')
        def _cd_entries(v):
            if isinstance(v, (int, float)):
                return [{'year': None, 'amount': v}] if v else []
            return v if isinstance(v, list) else []
        cd_str = ''
        if r.get('cross_dept'):
            parts = []
            for k, v in cd.items():
                entries = [e for e in _cd_entries(v) if e.get('amount')]
                if not entries:
                    continue
                detail = '、'.join(
                    f"{e['year']}年:{e['amount']:,.0f}" if e.get('year') else f"{e['amount']:,.0f}"
                    for e in entries
                )
                parts.append(f"{k}（{detail}）")
            cd_str = '；'.join(parts)
        vals = [r.get('client',''), r.get('project_name',''), r.get('group_name',''), r.get('status',''),
                r.get('expected_amount',0) or '', r.get('expected_date',''),
                r.get('amount',0) or '', r.get('sign_date',''),
                r.get('payment_type','當年'), r.get('installments',1) if r.get('payment_type')=='分期' else '',
                cd_str, '是' if r.get('carry_next') else '', r.get('note','')]
        for c,v in enumerate(vals,1):
            ws.cell(ri,c,v).border=border

    output = io.BytesIO(); wb.save(output); output.seek(0)
    return send_file(output, as_attachment=True,
                     download_name=f'{year}年{dept}{month}月合約.xlsx',
                     mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')

@app.route('/api/delete_contract/<int:cid>', methods=['DELETE'])
@login_required
def delete_contract(cid):
    con = get_db()
    row = con.execute('SELECT dept FROM contracts WHERE id=?', (cid,)).fetchone()
    if row and not can_write_dept(row['dept']):
        con.close()
        return jsonify({'error': '無寫入權限'}), 403
    con.execute('DELETE FROM contracts WHERE id=?', (cid,))
    con.commit()
    con.close()
    return jsonify({'status': 'ok'})

@app.route('/api/overview')
@login_required
def api_overview():
    year = get_current_year()
    con = get_db()
    result = []
    for dept in DEPARTMENTS:
        income   = con.execute("SELECT COALESCE(SUM(amount),0) FROM revenue WHERE year=? AND dept=? AND item='來自民間收入'", (year, dept)).fetchone()[0]
        expense  = con.execute("SELECT COALESCE(SUM(amount),0) FROM revenue WHERE year=? AND dept=? AND item='其他民間收入支出'", (year, dept)).fetchone()[0]
        unclaim  = con.execute("SELECT COALESCE(SUM(amount),0) FROM unclaimed WHERE year=? AND dept=?", (year, dept)).fetchone()[0]
        contracts= con.execute("SELECT COUNT(*) FROM contracts WHERE year=? AND dept=?", (year, dept)).fetchone()[0]
        result.append({'dept': dept, 'income': income, 'expense': expense + unclaim,
                       'unclaim': unclaim, 'contracts': contracts})
    con.close()
    return jsonify({'depts': result})

@app.route('/api/summary_data')
@login_required
def api_summary_data():
    year = request.args.get('year', type=int) or get_current_year()
    con = get_db()
    rows = con.execute(
        "SELECT dept, month, item, amount FROM revenue WHERE year=? AND item IN ('來自民間收入','其他民間收入支出')",
        (year,)
    ).fetchall()
    ucl = con.execute(
        "SELECT dept, month, SUM(amount) as total FROM unclaimed WHERE year=? GROUP BY dept, month",
        (year,)
    ).fetchall()
    goals = con.execute(
        "SELECT dept, item, goal FROM annual_goals WHERE year=? AND item IN ('來自民間收入','其他民間收入支出')",
        (year,)
    ).fetchall()
    con.close()
    goal_map = {(r['dept'], r['item']): r['goal'] for r in goals}
    data = {}
    for r in rows:
        d = data.setdefault(r['dept'], {}).setdefault(r['month'], {})
        if r['item'] == '來自民間收入':
            d['income'] = r['amount']
            d['income_goal'] = goal_map.get((r['dept'], '來自民間收入'), 0)
        else:
            d['expense'] = r['amount']
            d['expense_goal'] = goal_map.get((r['dept'], '其他民間收入支出'), 0)
    for r in ucl:
        data.setdefault(r['dept'], {}).setdefault(r['month'], {})['unclaim'] = r['total']
    return jsonify(data)

@app.route('/api/full_report')
@login_required
def api_full_report():
    year = request.args.get('year', type=int) or get_current_year()
    thru = request.args.get('thru_month', type=int) or 12
    con = get_db()
    ph = '%s' if IS_PG else '?'

    goals_rows = con.execute(
        f'SELECT item, SUM(goal) as total FROM annual_goals WHERE year={ph} GROUP BY item', (year,)
    ).fetchall()
    goals = {r['item']: (r['total'] or 0) for r in goals_rows}

    rev_rows = con.execute(
        f'SELECT item, month, SUM(amount) as amt, SUM(expected_amount) as exp '
        f'FROM revenue WHERE year={ph} GROUP BY item, month', (year,)
    ).fetchall()
    rev_by_item_month = {}
    for r in rev_rows:
        rev_by_item_month.setdefault(r['item'], {})[r['month']] = {
            'amt': r['amt'] or 0, 'exp': r['exp'] or 0
        }

    UNCL_MAP = {'業務費(未核銷)':'業務費','旅運費(未核銷)':'旅運費',
                '材料費(未核銷)':'材料費','維護費(未核銷)':'維護費'}
    ucl_rows = con.execute(
        f'SELECT item, month, SUM(amount) as amt FROM unclaimed WHERE year={ph} GROUP BY item, month', (year,)
    ).fetchall()
    ucl_by_exp_month = {}
    for r in ucl_rows:
        exp_item = UNCL_MAP.get(r['item'])
        if exp_item:
            ucl_by_exp_month.setdefault(exp_item, {})[r['month']] = \
                ucl_by_exp_month.get(exp_item, {}).get(r['month'], 0) + (r['amt'] or 0)
    con.close()

    INCOME_LABELS = {'其他民間收入': '其他民間收入(試驗/技術/訓練/其他)'}

    def ytd(item_map, m):
        return sum((item_map.get(mm, {}).get('amt', 0) for mm in range(1, m+1)), 0)
    def ytd_exp(item_map, m):
        return sum((item_map.get(mm, {}).get('exp', 0) for mm in range(1, m+1)), 0)
    def ytd_ucl(ucl_map, m):
        return sum((ucl_map.get(mm, 0) for mm in range(1, m+1)), 0)

    income = []
    for item in INCOME_ITEMS:
        im = rev_by_item_month.get(item, {})
        g = goals.get(item, 0)
        signed = ytd(im, thru)
        expected = ytd_exp(im, thru)
        monthly = {m: im.get(m, {}).get('amt', 0) for m in range(1,13)}
        monthly_exp = {m: im.get(m, {}).get('exp', 0) for m in range(1,13)}
        income.append({
            'item': item, 'label': INCOME_LABELS.get(item, item),
            'goal': g, 'signed': signed, 'expected': expected,
            'subtotal': signed + expected, 'diff': (signed + expected) - g,
            'monthly': monthly, 'monthly_exp': monthly_exp,
            'is_total': item == INCOME_TOTAL_ITEM
        })

    expense = []
    for item in EXPENSE_ITEMS:
        im = rev_by_item_month.get(item, {})
        um = ucl_by_exp_month.get(item, {})
        g = goals.get(item, 0)
        actual = ytd(im, thru)
        unclaimed = ytd_ucl(um, thru)
        monthly = {m: im.get(m, {}).get('amt', 0) for m in range(1,13)}
        monthly_ucl = {m: um.get(m, 0) for m in range(1,13)}
        expense.append({
            'item': item, 'goal': g, 'actual': actual, 'unclaimed': unclaimed,
            'subtotal': actual + unclaimed, 'diff': (actual + unclaimed) - g,
            'monthly': monthly, 'monthly_unclaimed': monthly_ucl,
            'is_total': item == '其他民間收入支出'
        })

    return jsonify({'income': income, 'expense': expense,
                    'year': year, 'thru_month': thru})

@app.template_filter('enumerate')
def enumerate_filter(iterable, start=0):
    return list(enumerate(iterable, start))

@app.route('/summary')
@login_required
def summary():
    year = get_current_year()
    all_years = get_all_years()
    con = get_db()
    rows = con.execute('''
        SELECT dept, month, item, SUM(amount) as total
        FROM revenue WHERE year=? AND item IN ('來自民間業務收入合計', '支出合計')
        GROUP BY dept, month, item
    ''', (year,)).fetchall()
    con.close()
    data = {}
    for r in rows:
        data.setdefault(r['dept'], {}).setdefault(r['month'], {})[r['item']] = r['total']
    allowed = get_allowed_depts()
    return render_template('summary.html', departments=allowed, months=MONTHS,
                           year=year, all_years=all_years, data=data)

@app.route('/contracts')
@login_required
def contracts_view():
    year  = get_current_year()
    all_years = get_all_years()
    dept  = request.args.get('dept', '')
    month = request.args.get('month', 0, type=int)
    con = get_db()
    q, params = 'SELECT * FROM contracts WHERE year=?', [year]
    if dept:  q += ' AND dept=?';  params.append(dept)
    if month: q += ' AND month=?'; params.append(month)
    q += ' ORDER BY month, dept, id'
    contracts = [dict(r) for r in con.execute(q, params).fetchall()]
    con.close()
    allowed = get_allowed_depts()
    return render_template('contracts.html', contracts=contracts,
                           departments=allowed, months=MONTHS,
                           year=year, all_years=all_years,
                           selected_dept=dept, selected_month=month,
                           statuses=CONTRACT_STATUSES)

@app.route('/dept/<dept>/contracts')
@login_required
def dept_contracts(dept):
    if dept not in DEPARTMENTS:
        return redirect(url_for('index'))
    if not can_access_dept(dept):
        return redirect(url_for('dept_contracts', dept=session.get('dept'),
                                month=request.args.get('month', 1)))
    month = request.args.get('month', 1, type=int)
    year = get_current_year()
    all_years = get_all_years()
    allowed = get_allowed_depts()
    return render_template('dept_contracts.html', dept=dept, month=month,
                           months=MONTHS, year=year, all_years=all_years,
                           departments=allowed,
                           can_write=can_write_dept(dept))

@app.route('/import')
@login_required
def import_view():
    year = get_current_year()
    all_years = get_all_years()
    return render_template('import.html', departments=DEPARTMENTS, months=MONTHS,
                           year=year, all_years=all_years)

@app.route('/download_import_template')
@login_required
def download_import_template():
    year = get_current_year()
    wb = openpyxl.Workbook()
    thin   = Side(style='thin')
    border = Border(left=thin, right=thin, top=thin, bottom=thin)
    hfill  = PatternFill('solid', start_color='BDD7EE', fgColor='BDD7EE')
    hfont  = Font(name='微軟正黑體', size=10, bold=True)
    ctr    = Alignment(horizontal='center', vertical='center')
    note_fill = PatternFill('solid', start_color='FFFF99', fgColor='FFFF99')

    def make_hdr(ws, headers, widths):
        for col, (h, w) in enumerate(zip(headers, widths), 1):
            c = ws.cell(row=1, column=col, value=h)
            c.font = hfont; c.fill = hfill; c.alignment = ctr; c.border = border
            ws.column_dimensions[get_column_letter(col)].width = w

    # Sheet1: 收支資料
    ws1 = wb.active; ws1.title = '收支資料'
    make_hdr(ws1, ['部門','月份','項目','金額','年度目標'],
             [12, 8, 28, 14, 14])
    for ri, (dept, item) in enumerate(
        [(d, i) for d in DEPARTMENTS for i in (INCOME_ITEMS + EXPENSE_ITEMS)], 2):
        ws1.cell(row=ri, column=1, value=dept).border = border
        ws1.cell(row=ri, column=2, value=1).border = border
        ws1.cell(row=ri, column=3, value=item).border = border
        ws1.cell(row=ri, column=4, value=0).border = border
        ws1.cell(row=ri, column=5, value=0).border = border
    note = ws1.cell(row=1, column=7, value='※ 請勿修改「項目」欄位文字；月份填 1-12；部門名稱需完全符合')
    note.fill = note_fill; note.font = Font(name='微軟正黑體', size=9, color='FF0000')

    # Sheet2: 未核銷費用
    ws2 = wb.create_sheet('未核銷費用')
    make_hdr(ws2, ['部門','月份','未核銷項目','金額'],
             [12, 8, 22, 14])
    for ri, (dept, item) in enumerate(
        [(d, i) for d in DEPARTMENTS for i in UNCLAIMED_ITEMS], 2):
        ws2.cell(row=ri, column=1, value=dept).border = border
        ws2.cell(row=ri, column=2, value=1).border = border
        ws2.cell(row=ri, column=3, value=item).border = border
        ws2.cell(row=ri, column=4, value=0).border = border

    # Sheet3: 合約
    ws3 = wb.create_sheet('合約追蹤')
    make_hdr(ws3, ['部門','月份','客戶/計畫名稱','合約金額','簽約日期','預計完成日','狀態','本月實收金額','備註','延續下月(是/否)'],
             [12, 8, 28, 14, 14, 14, 14, 14, 20, 12])
    for ri, dept in enumerate(DEPARTMENTS, 2):
        ws3.cell(row=ri, column=1, value=dept).border = border
        ws3.cell(row=ri, column=2, value=1).border = border
        for c in range(3, 11):
            ws3.cell(row=ri, column=c, value='').border = border
    note3 = ws3.cell(row=1, column=12, value='※ 狀態可填：洽談中、新增簽約、已簽約執行中、完成')
    note3.fill = note_fill; note3.font = Font(name='微軟正黑體', size=9, color='FF0000')

    output = io.BytesIO()
    wb.save(output); output.seek(0)
    return send_file(output, as_attachment=True,
                     download_name=f'{year}年收支資料匯入範本.xlsx',
                     mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')

@app.route('/api/import_excel', methods=['POST'])
@login_required
def import_excel():
    if not can_write():
        return jsonify({'error': '無寫入權限'}), 403
    f = request.files.get('file')
    if not f:
        return jsonify({'error': '未收到檔案'}), 400
    year = get_current_year()
    try:
        wb = openpyxl.load_workbook(f, data_only=True)
    except Exception as e:
        return jsonify({'error': f'檔案格式錯誤：{e}'}), 400

    con = get_db()
    stats = {'revenue': 0, 'unclaimed': 0, 'contracts': 0, 'errors': []}

    # Sheet: 收支資料
    if '收支資料' in wb.sheetnames:
        ws = wb['收支資料']
        headers = [ws.cell(row=1, column=c).value for c in range(1, 6)]
        for row in ws.iter_rows(min_row=2, values_only=True):
            if not any(row): continue
            dept, month, item, amount, goal = (row + (None,)*5)[:5]
            if dept not in DEPARTMENTS:
                stats['errors'].append(f'收支資料：部門 {dept!r} 不存在，略過')
                continue
            if item not in (INCOME_ITEMS + EXPENSE_ITEMS):
                stats['errors'].append(f'收支資料：項目 {item!r} 不存在，略過')
                continue
            try:
                month = int(month or 1)
                amount = float(amount or 0)
                goal = float(goal or 0)
            except:
                stats['errors'].append(f'收支資料：{dept}/{item} 數值格式錯誤，略過')
                continue
            con.execute('''INSERT INTO revenue (year, dept, month, item, amount, goal)
                VALUES (?,?,?,?,?,?)
                ON CONFLICT(year, dept, month, item) DO UPDATE
                SET amount=excluded.amount, goal=excluded.goal, updated_at=CURRENT_TIMESTAMP''',
                (year, dept, month, item, amount, goal))
            stats['revenue'] += 1

    # Sheet: 未核銷費用
    if '未核銷費用' in wb.sheetnames:
        ws = wb['未核銷費用']
        for row in ws.iter_rows(min_row=2, values_only=True):
            if not any(row): continue
            dept, month, item, amount = (row + (None,)*4)[:4]
            if dept not in DEPARTMENTS:
                stats['errors'].append(f'未核銷：部門 {dept!r} 不存在，略過')
                continue
            if item not in UNCLAIMED_ITEMS:
                stats['errors'].append(f'未核銷：項目 {item!r} 不存在，略過')
                continue
            try:
                month = int(month or 1); amount = float(amount or 0)
            except:
                continue
            con.execute('''INSERT INTO unclaimed (year, dept, month, item, amount)
                VALUES (?,?,?,?,?)
                ON CONFLICT(year, dept, month, item) DO UPDATE
                SET amount=excluded.amount, updated_at=CURRENT_TIMESTAMP''',
                (year, dept, month, item, amount))
            stats['unclaimed'] += 1

    # Sheet: 合約追蹤
    if '合約追蹤' in wb.sheetnames:
        ws = wb['合約追蹤']
        for row in ws.iter_rows(min_row=2, values_only=True):
            if not any(row): continue
            dept, month, client, amount, sign_date, due_date, status, actual, note, carry = (row + (None,)*10)[:10]
            if dept not in DEPARTMENTS:
                stats['errors'].append(f'合約：部門 {dept!r} 不存在，略過')
                continue
            if not client:
                continue
            try:
                month = int(month or 1)
                amount = float(amount or 0)
                actual = float(actual or 0)
            except:
                continue
            carry_val = 1 if str(carry or '').strip() in ('是', '1', 'True', 'true', 'Y', 'y') else 0
            status = str(status or '洽談中').strip()
            if status not in CONTRACT_STATUSES:
                status = '洽談中'
            con.execute('''INSERT INTO contracts
                (year, dept, month, client, amount, sign_date, due_date,
                 status, actual_amount, note, carry_next)
                VALUES (?,?,?,?,?,?,?,?,?,?,?)''',
                (year, dept, month, str(client), amount,
                 str(sign_date or ''), str(due_date or ''),
                 status, actual, str(note or ''), carry_val))
            stats['contracts'] += 1

    con.commit(); con.close()
    msg = f"匯入完成：收支資料 {stats['revenue']} 筆，未核銷 {stats['unclaimed']} 筆，合約 {stats['contracts']} 筆"
    if stats['errors']:
        msg += f"；{len(stats['errors'])} 筆略過（見詳情）"
    return jsonify({'status': 'ok', 'message': msg, 'errors': stats['errors'][:20]})

@app.route('/export_pptx')
@login_required
def export_pptx():
    import traceback
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.oxml.ns import qn
        import copy
    except ImportError as e:
        return f'<h2>缺少套件</h2><pre>{e}\n\n請在 PythonAnywhere 執行:\npip install python-pptx</pre>', 500

    year = request.args.get('year', type=int) or get_current_year()
    thru = request.args.get('thru_month', type=int) or 12
    ph = '%s' if IS_PG else '?'
    con = get_db()

    # YTD 收入/支出資料
    rev_rows = con.execute(
        f'SELECT dept, item, SUM(amount) as amt, SUM(expected_amount) as exp '
        f'FROM revenue WHERE year={ph} AND month<={ph} GROUP BY dept, item',
        (year, thru)
    ).fetchall()
    uncl_rows = con.execute(
        f'SELECT dept, item, SUM(amount) as amt FROM unclaimed '
        f'WHERE year={ph} AND month<={ph} GROUP BY dept, item',
        (year, thru)
    ).fetchall()
    goals_rows = con.execute(
        f'SELECT dept, item, goal FROM annual_goals WHERE year={ph}', (year,)
    ).fetchall()
    contract_rows = con.execute(
        f'SELECT * FROM contracts WHERE year={ph} ORDER BY dept, status', (year,)
    ).fetchall()
    con.close()

    # 建立資料結構
    cum_rev = {}   # {dept: {item: {amt, exp}}}
    for r in rev_rows:
        cum_rev.setdefault(r['dept'], {})[r['item']] = {'amt': r['amt'] or 0, 'exp': r['exp'] or 0}

    cum_uncl = {}  # {dept: {item: amount}}
    UNCL_MAP_R = {'業務費(未核銷)':'業務費','旅運費(未核銷)':'旅運費',
                  '材料費(未核銷)':'材料費','維護費(未核銷)':'維護費'}
    for r in uncl_rows:
        exp_item = UNCL_MAP_R.get(r['item'], r['item'])
        cum_uncl.setdefault(r['dept'], {})[exp_item] = \
            cum_uncl.get(r['dept'], {}).get(exp_item, 0) + (r['amt'] or 0)

    goals = {}  # {dept: {item: goal}}
    for r in goals_rows:
        goals.setdefault(r['dept'], {})[r['item']] = r['goal'] or 0

    INCOME_LABEL = {'其他民間收入': '其他民間收入(試驗/技術/訓練/其他)'}
    DEPT_NUM = {'原料部':'一','產品部':'二','檢驗部':'三','製程部':'四','雲分部':'五','產服部':'六'}
    INCOME_DISPLAY = [i for i in INCOME_ITEMS if i != INCOME_TOTAL_ITEM]
    EXPENSE_DISPLAY = [i for i in EXPENSE_ITEMS if i != '其他民間收入支出']

    def fmtv(v):
        if v is None or v == 0: return '-'
        return f'{v:,.0f}'
    def diff_str(a, g):
        if not g: return '-'
        d = a - g
        return f'{d:+,.0f}'
    def rate_str(a, g):
        if not g: return '-'
        return f'{a/g*100:.1f}%'

    # 建立全所合計資料（不分部門）
    def agg_all(item):
        amt = sum(cum_rev.get(d, {}).get(item, {}).get('amt', 0) for d in DEPARTMENTS)
        exp_v = sum(cum_rev.get(d, {}).get(item, {}).get('exp', 0) for d in DEPARTMENTS)
        uncl = sum(cum_uncl.get(d, {}).get(item, 0) for d in DEPARTMENTS)
        goal = sum(goals.get(d, {}).get(item, 0) for d in DEPARTMENTS)
        return amt, exp_v, uncl, goal

    DEPT_NAMES = {d: n for d, n in DEPT_NUM.items()}

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # 全空白

    # ── 顏色常數 ──────────────────────────
    C_DARK_BLUE  = RGBColor(0x1F, 0x4E, 0x79)
    C_MID_BLUE   = RGBColor(0x2E, 0x75, 0xB6)
    C_INC_LIGHT  = RGBColor(0xE2, 0xEF, 0xDA)  # 收入淺綠
    C_INC_TOTAL  = RGBColor(0xA9, 0xD1, 0x8E)  # 收入合計
    C_EXP_LIGHT  = RGBColor(0xFF, 0xE0, 0xE0)  # 支出淺紅
    C_EXP_TOTAL  = RGBColor(0xFF, 0xB0, 0xB0)  # 支出合計
    C_RATE_BG    = RGBColor(0xFF, 0xFF, 0xCC)  # 達成率
    C_NET_BG     = RGBColor(0xFF, 0xD9, 0x66)  # 業務餘絀
    C_HDR_INC    = RGBColor(0x37, 0x86, 0x10)  # 收入表頭文字
    C_HDR_EXP    = RGBColor(0xC0, 0x00, 0x00)  # 支出表頭文字
    C_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)

    def set_cell(cell, text, bold=False, size=9, bg=None, color=None, align='center', wrap=True):
        tf = cell.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.alignment = {'center': PP_ALIGN.CENTER, 'left': PP_ALIGN.LEFT, 'right': PP_ALIGN.RIGHT}[align]
        run = p.add_run()
        run.text = str(text)
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = '微軟正黑體'
        if color:
            run.font.color.rgb = color
        if bg:
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            solidFill = tc.makeelement(qn('a:solidFill'))
            srgbClr = tc.makeelement(qn('a:srgbClr'))
            srgbClr.set('val', f'{bg.red:02X}{bg.green:02X}{bg.blue:02X}')
            solidFill.append(srgbClr)
            tcPr.append(solidFill)

    # ── 封面投影片 ──────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    # 背景色
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = C_DARK_BLUE

    tx = slide.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(10.0), Inches(1.5))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f'業務部門1-{thru}月來自民間業務收入'
    run.font.size = Pt(32); run.font.bold = True; run.font.color.rgb = C_WHITE; run.font.name = '微軟正黑體'

    tx2 = slide.shapes.add_textbox(Inches(1.5), Inches(3.5), Inches(10.0), Inches(0.8))
    tf2 = tx2.text_frame; p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER; run2 = p2.add_run()
    run2.text = f'中華民國{year}年{thru}月'
    run2.font.size = Pt(20); run2.font.color.rgb = RGBColor(0xBD, 0xD7, 0xEE); run2.font.name = '微軟正黑體'

    tx3 = slide.shapes.add_textbox(Inches(1.5), Inches(4.2), Inches(10.0), Inches(0.6))
    tf3 = tx3.text_frame; p3 = tf3.paragraphs[0]; p3.alignment = PP_ALIGN.CENTER; run3 = p3.add_run()
    run3.text = f'紡織所 企劃處'
    run3.font.size = Pt(14); run3.font.color.rgb = RGBColor(0xBD, 0xD7, 0xEE); run3.font.name = '微軟正黑體'

    # ── 共用：產生一張部門收支投影片 ──────────────────────
    def make_dept_slide(dept_label, dept_rev_d, dept_uncl_d, dept_goals_d, page_no):
        slide = prs.slides.add_slide(blank_layout)

        # 部門標題框
        title_box = slide.shapes.add_textbox(Inches(0.2), Inches(0.1), Inches(5.0), Inches(0.45))
        tf = title_box.text_frame
        # 部門標題
        p = tf.paragraphs[0]; run = p.add_run()
        run.text = dept_label
        run.font.size = Pt(16); run.font.bold = True; run.font.color.rgb = C_DARK_BLUE; run.font.name = '微軟正黑體'

        unit_box = slide.shapes.add_textbox(Inches(11.8), Inches(0.05), Inches(1.3), Inches(0.35))
        tf_u = unit_box.text_frame; p_u = tf_u.paragraphs[0]; p_u.alignment = PP_ALIGN.RIGHT
        run_u = p_u.add_run(); run_u.text = '單位：元'
        run_u.font.size = Pt(8); run_u.font.color.rgb = RGBColor(0x60, 0x60, 0x60); run_u.font.name = '微軟正黑體'

        pg_box = slide.shapes.add_textbox(Inches(12.8), Inches(7.1), Inches(0.4), Inches(0.3))
        tf_p = pg_box.text_frame; p_pg = tf_p.paragraphs[0]; p_pg.alignment = PP_ALIGN.RIGHT
        run_p = p_pg.add_run(); run_p.text = str(page_no)
        run_p.font.size = Pt(9); run_p.font.color.rgb = RGBColor(0x80, 0x80, 0x80); run_p.font.name = '微軟正黑體'

        # 篩選有資料的收入/支出行
        inc_rows = [i for i in INCOME_DISPLAY
                    if dept_rev_d.get(i,{}).get('amt',0) or dept_rev_d.get(i,{}).get('exp',0) or dept_goals_d.get(i,0)]
        exp_rows = [i for i in EXPENSE_DISPLAY
                    if dept_rev_d.get(i,{}).get('amt',0) or dept_uncl_d.get(i,0) or dept_goals_d.get(i,0)]

        # 總收入/支出
        x_signed  = dept_rev_d.get(INCOME_TOTAL_ITEM,{}).get('amt',0) or sum(dept_rev_d.get(i,{}).get('amt',0) for i in inc_rows)
        x_exp     = dept_rev_d.get(INCOME_TOTAL_ITEM,{}).get('exp',0) or sum(dept_rev_d.get(i,{}).get('exp',0) for i in inc_rows)
        x_goal    = dept_goals_d.get(INCOME_TOTAL_ITEM,0) or sum(dept_goals_d.get(i,0) for i in inc_rows)
        x_subtot  = x_signed + x_exp

        exp_actual = dept_rev_d.get('其他民間收入支出',{}).get('amt',0) or sum(dept_rev_d.get(i,{}).get('amt',0) for i in exp_rows)
        exp_uncl   = sum(dept_uncl_d.get(i,0) for i in exp_rows)
        exp_goal   = dept_goals_d.get('其他民間收入支出',0) or sum(dept_goals_d.get(i,0) for i in exp_rows)
        exp_subtot = exp_actual + exp_uncl

        svc_signed = dept_rev_d.get('其他民間收入',{}).get('amt',0)
        svc_exp    = dept_rev_d.get('其他民間收入',{}).get('exp',0)
        mgmt_fee   = round((svc_signed + svc_exp) * 0.115)
        y_val      = x_subtot - exp_subtot - mgmt_fee
        y_goal     = dept_goals_d.get('財務貢獻(Y值)', 0)

        # 表格結構：1行標頭 + inc_rows + 3行(X合計/達成率) + 1行支出頭 + exp_rows + 5行(合計/餘絀/管費/Y值/Y率)
        INC_R  = 1 + len(inc_rows) + 2
        EXP_R  = 1 + len(exp_rows) + 5
        NROWS  = INC_R + EXP_R
        COLS   = 7
        COL_W  = [Inches(2.85), Inches(1.42), Inches(1.42), Inches(1.42), Inches(1.42), Inches(1.42), Inches(1.42)]
        tbl_h  = Inches(6.85)
        table  = slide.shapes.add_table(NROWS, COLS, Inches(0.15), Inches(0.52), sum(COL_W), tbl_h).table
        for ci, w in enumerate(COL_W): table.columns[ci].width = w
        row_h = int(tbl_h / NROWS)
        for ri in range(NROWS): table.rows[ri].height = row_h

        def sc(ri, ci, txt, bold=False, bg=None, color=None, align='center', sz=8):
            set_cell(table.cell(ri,ci), txt, bold=bold, size=sz, bg=bg, color=color, align=align)

        C_RED = RGBColor(0xC0,0x00,0x00)
        C_INC_H = RGBColor(0x37,0x86,0x10)

        # 收入標頭
        sc(0,0,'收入項目',bold=True,bg=C_INC_H,color=C_WHITE,align='left')
        sc(0,1,'年度預算(A)',bold=True,bg=C_INC_H,color=C_WHITE)
        sc(0,2,f'1-{thru}月 已簽約/已收',bold=True,bg=C_INC_H,color=C_WHITE)
        sc(0,3,'預計簽約',bold=True,bg=C_INC_H,color=C_WHITE)
        sc(0,4,'小計(B)',bold=True,bg=C_INC_H,color=C_WHITE)
        sc(0,5,'差異=(B)-(A)',bold=True,bg=C_INC_H,color=C_WHITE)
        sc(0,6,'年度預估',bold=True,bg=C_INC_H,color=C_WHITE)

        # 收入資料行
        for i, item in enumerate(inc_rows):
            label = INCOME_LABEL.get(item, item)
            d = dept_rev_d.get(item, {})
            signed = d.get('amt', 0); exp_v = d.get('exp', 0)
            goal = dept_goals_d.get(item, 0)
            subtot = signed + exp_v
            ri = 1 + i
            sc(ri,0,label,bg=C_INC_LIGHT,align='left')
            sc(ri,1,fmtv(goal),bg=C_INC_LIGHT,align='right')
            sc(ri,2,fmtv(signed),bg=C_INC_LIGHT,align='right')
            sc(ri,3,fmtv(exp_v),bg=C_INC_LIGHT,align='right')
            sc(ri,4,fmtv(subtot),bg=C_INC_LIGHT,align='right')
            sc(ri,5,diff_str(subtot,goal),bg=C_INC_LIGHT,align='right',
               color=C_RED if subtot<goal else RGBColor(0x1A,0x7A,0x3C))
            sc(ri,6,fmtv(goal),bg=C_INC_LIGHT,align='right')

        # X值合計
        ri_x = 1 + len(inc_rows)
        sc(ri_x,0,'來自民間收入合計(X值)(A)',bold=True,bg=C_INC_TOTAL,align='left')
        sc(ri_x,1,fmtv(x_goal),bold=True,bg=C_INC_TOTAL,align='right')
        sc(ri_x,2,fmtv(x_signed),bold=True,bg=C_INC_TOTAL,align='right')
        sc(ri_x,3,fmtv(x_exp),bold=True,bg=C_INC_TOTAL,align='right')
        sc(ri_x,4,fmtv(x_subtot),bold=True,bg=C_INC_TOTAL,align='right')
        sc(ri_x,5,diff_str(x_subtot,x_goal),bold=True,bg=C_INC_TOTAL,align='right',
           color=C_RED if x_subtot<x_goal else RGBColor(0x1A,0x7A,0x3C))
        sc(ri_x,6,fmtv(x_goal),bold=True,bg=C_INC_TOTAL,align='right')

        # X值達成率
        ri_xr = ri_x + 1
        sc(ri_xr,0,'X值達成率',bold=True,bg=C_RATE_BG,align='left')
        sc(ri_xr,1,'',bg=C_RATE_BG)
        sc(ri_xr,2,rate_str(x_signed,x_goal),bold=True,bg=C_RATE_BG)
        sc(ri_xr,3,rate_str(x_exp,x_goal) if x_exp else '-',bg=C_RATE_BG)
        sc(ri_xr,4,rate_str(x_subtot,x_goal),bold=True,bg=C_RATE_BG)
        for ci in [5,6]: sc(ri_xr,ci,'',bg=C_RATE_BG)

        # 支出標頭
        ri_eh = INC_R
        sc(ri_eh,0,'費用項目',bold=True,bg=C_RED,color=C_WHITE,align='left')
        sc(ri_eh,1,'年度預算(1)',bold=True,bg=C_RED,color=C_WHITE)
        sc(ri_eh,2,'實際報支',bold=True,bg=C_RED,color=C_WHITE)
        sc(ri_eh,3,'已申請未銷核',bold=True,bg=C_RED,color=C_WHITE)
        sc(ri_eh,4,'小計(2)',bold=True,bg=C_RED,color=C_WHITE)
        sc(ri_eh,5,'差異=(2)-(1)',bold=True,bg=C_RED,color=C_WHITE)
        sc(ri_eh,6,'年度預估',bold=True,bg=C_RED,color=C_WHITE)

        # 支出資料行
        for i, item in enumerate(exp_rows):
            d = dept_rev_d.get(item, {})
            actual = d.get('amt', 0)
            uncl = dept_uncl_d.get(item, 0)
            goal = dept_goals_d.get(item, 0)
            subtot = actual + uncl
            ri = ri_eh + 1 + i
            sc(ri,0,item,bg=C_EXP_LIGHT,align='left')
            sc(ri,1,fmtv(goal),bg=C_EXP_LIGHT,align='right')
            sc(ri,2,fmtv(actual),bg=C_EXP_LIGHT,align='right')
            sc(ri,3,fmtv(uncl) if uncl else '-',bg=C_EXP_LIGHT,align='right')
            sc(ri,4,fmtv(subtot),bg=C_EXP_LIGHT,align='right')
            sc(ri,5,diff_str(subtot,goal),bg=C_EXP_LIGHT,align='right')
            sc(ri,6,fmtv(goal),bg=C_EXP_LIGHT,align='right')

        # 支出合計
        ri_et = ri_eh + 1 + len(exp_rows)
        sc(ri_et,0,'其他民間收入支出',bold=True,bg=C_EXP_TOTAL,align='left')
        sc(ri_et,1,fmtv(exp_goal),bold=True,bg=C_EXP_TOTAL,align='right')
        sc(ri_et,2,fmtv(exp_actual),bold=True,bg=C_EXP_TOTAL,align='right')
        sc(ri_et,3,fmtv(exp_uncl),bold=True,bg=C_EXP_TOTAL,align='right')
        sc(ri_et,4,fmtv(exp_subtot),bold=True,bg=C_EXP_TOTAL,align='right')
        sc(ri_et,5,diff_str(exp_subtot,exp_goal),bold=True,bg=C_EXP_TOTAL,align='right')
        sc(ri_et,6,fmtv(exp_goal),bold=True,bg=C_EXP_TOTAL,align='right')

        # 業務餘絀
        ri_net = ri_et + 1; net_val = x_subtot - exp_subtot
        sc(ri_net,0,'業務餘絀',bold=True,bg=C_NET_BG,align='left')
        sc(ri_net,1,'-',bg=C_NET_BG); sc(ri_net,2,fmtv(x_signed-exp_actual),bg=C_NET_BG,align='right')
        sc(ri_net,3,'-',bg=C_NET_BG); sc(ri_net,4,fmtv(net_val),bold=True,bg=C_NET_BG,align='right',
           color=C_DARK_BLUE if net_val>=0 else C_RED)
        sc(ri_net,5,'-',bg=C_NET_BG); sc(ri_net,6,'-',bg=C_NET_BG)

        # 管理費
        ri_mg = ri_net + 1; svc_g = dept_goals_d.get('其他民間收入',0)
        sc(ri_mg,0,'管理費(×11.5%)',bold=True,bg=C_RATE_BG,align='left')
        sc(ri_mg,1,fmtv(round(svc_g*0.115)),bg=C_RATE_BG,align='right')
        sc(ri_mg,2,fmtv(round(svc_signed*0.115)),bg=C_RATE_BG,align='right')
        sc(ri_mg,3,'-',bg=C_RATE_BG); sc(ri_mg,4,fmtv(mgmt_fee),bg=C_RATE_BG,align='right')
        sc(ri_mg,5,'-',bg=C_RATE_BG); sc(ri_mg,6,'-',bg=C_RATE_BG)

        # 財務貢獻Y值
        ri_y = ri_mg + 1
        C_YBG = RGBColor(0xDC,0xE6,0xF1)
        sc(ri_y,0,'財務貢獻(Y值)',bold=True,bg=C_YBG,align='left',color=C_DARK_BLUE)
        sc(ri_y,1,fmtv(y_goal),bold=True,bg=C_YBG,align='right',color=C_DARK_BLUE)
        sc(ri_y,2,fmtv(x_signed-exp_actual-round(svc_signed*0.115)),bold=True,bg=C_YBG,align='right',color=C_DARK_BLUE)
        sc(ri_y,3,'-',bg=C_YBG); sc(ri_y,4,fmtv(y_val),bold=True,bg=C_YBG,align='right',
           color=C_DARK_BLUE if y_val>=0 else C_RED)
        sc(ri_y,5,diff_str(y_val,y_goal),bold=True,bg=C_YBG,align='right',
           color=C_RED if y_val<y_goal else RGBColor(0x1A,0x7A,0x3C))
        sc(ri_y,6,'-',bg=C_YBG)

        # Y值達成率
        ri_yr = ri_y + 1
        sc(ri_yr,0,'Y值達成率',bold=True,bg=C_RATE_BG,align='left')
        sc(ri_yr,1,'',bg=C_RATE_BG)
        sc(ri_yr,2,rate_str(x_signed-exp_actual-round(svc_signed*0.115), y_goal),bold=True,bg=C_RATE_BG)
        sc(ri_yr,3,'-',bg=C_RATE_BG)
        sc(ri_yr,4,rate_str(y_val,y_goal),bold=True,bg=C_RATE_BG)
        for ci in [5,6]: sc(ri_yr,ci,'',bg=C_RATE_BG)

    # ── 各部門投影片 ─────────────────────────────────────
    for dept_idx, dept in enumerate(DEPARTMENTS):
        dept_num = DEPT_NAMES[dept]
        make_dept_slide(f'{dept_num}、{dept}',
                        cum_rev.get(dept, {}), cum_uncl.get(dept, {}),
                        goals.get(dept, {}), dept_idx + 1)

    # ── 全所合計投影片 ──────────────────────────────────
    all_rev = {}
    all_uncl = {}
    all_goals = {}
    for d in DEPARTMENTS:
        for item, v in cum_rev.get(d, {}).items():
            r = all_rev.setdefault(item, {'amt':0,'exp':0})
            r['amt'] += v.get('amt',0); r['exp'] += v.get('exp',0)
        for item, v in cum_uncl.get(d, {}).items():
            all_uncl[item] = all_uncl.get(item,0) + v
        for item, v in goals.get(d, {}).items():
            all_goals[item] = all_goals.get(item,0) + v

    make_dept_slide(f'七、業務部門（全所合計）', all_rev, all_uncl, all_goals, len(DEPARTMENTS)+1)

    # ── 合約附件投影片 ──────────────────────────────────
    contracts_by_dept = {}
    for c in contract_rows:
        contracts_by_dept.setdefault(dict(c)['dept'], []).append(dict(c))

    DEPT_NUM_MAP = {d: n for n, d in enumerate(DEPARTMENTS, 1)}
    STATUS_ORDER = ['已簽約','簽約中/用印中','預計簽約','洽談中','停止']
    append_page = 1
    depts_per_slide = 3
    dept_chunks = [DEPARTMENTS[i:i+depts_per_slide] for i in range(0, len(DEPARTMENTS), depts_per_slide)]

    for chunk_idx, chunk in enumerate(dept_chunks):
        slide = prs.slides.add_slide(blank_layout)
        title_tx = slide.shapes.add_textbox(Inches(0.2), Inches(0.05), Inches(10.0), Inches(0.4))
        tf_t = title_tx.text_frame; p_t = tf_t.paragraphs[0]
        run_t = p_t.add_run()
        run_t.text = f'附件、來自民間業務簽約狀況({append_page}/{len(dept_chunks)})'
        run_t.font.size = Pt(12); run_t.font.bold = True; run_t.font.color.rgb = C_DARK_BLUE; run_t.font.name = '微軟正黑體'

        pg_bx = slide.shapes.add_textbox(Inches(12.8), Inches(7.1), Inches(0.4), Inches(0.3))
        tf_pg = pg_bx.text_frame; p_pg2 = tf_pg.paragraphs[0]; p_pg2.alignment = PP_ALIGN.RIGHT
        run_pg = p_pg2.add_run(); run_pg.text = str(len(DEPARTMENTS)+1+append_page)
        run_pg.font.size = Pt(9); run_pg.font.color.rgb = RGBColor(0x80,0x80,0x80); run_pg.font.name = '微軟正黑體'
        append_page += 1

        y_pos = Inches(0.5)
        for dept in chunk:
            dept_contracts = sorted(contracts_by_dept.get(dept, []),
                                    key=lambda c: STATUS_ORDER.index(c.get('status','洽談中')) if c.get('status') in STATUS_ORDER else 99)
            if not dept_contracts:
                continue
            dept_num_str = DEPT_NAMES.get(dept, '')

            # 部門小標題
            lbl_bx = slide.shapes.add_textbox(Inches(0.2), y_pos, Inches(4.0), Inches(0.3))
            tf_l = lbl_bx.text_frame; p_l = tf_l.paragraphs[0]
            run_l = p_l.add_run(); run_l.text = f'({dept_num_str}){dept}'
            run_l.font.size = Pt(11); run_l.font.bold = True; run_l.font.color.rgb = C_DARK_BLUE; run_l.font.name = '微軟正黑體'
            y_pos += Inches(0.32)

            # 合約表
            tbl_cols = 5
            n_rows = 1 + len(dept_contracts)
            tbl_h2 = Inches(0.22 * n_rows + 0.1)
            ct = slide.shapes.add_table(n_rows, tbl_cols, Inches(0.2), y_pos, Inches(13.0), tbl_h2).table
            cww = [Inches(0.5), Inches(3.5), Inches(3.5), Inches(1.5), Inches(4.0)]
            for ci2, w2 in enumerate(cww): ct.columns[ci2].width = w2

            for ci2, hdr_txt in enumerate(['編號','洽談廠商/客戶','計畫名稱','負責組別','狀況']):
                set_cell(ct.cell(0,ci2), hdr_txt, bold=True, size=8, bg=C_DARK_BLUE, color=C_WHITE)

            for ri2, c in enumerate(dept_contracts, 1):
                status = c.get('status','')
                status_bg = {
                    '已簽約': RGBColor(0xC6,0xEF,0xCE),
                    '簽約中/用印中': RGBColor(0xD0,0xE4,0xFF),
                    '預計簽約': RGBColor(0xFF,0xEB,0x9C),
                    '洽談中': C_WHITE, '停止': RGBColor(0xF2,0xDC,0xDB),
                }.get(status, C_WHITE)

                sign_info = ''
                if status in ('已簽約','簽約中/用印中') and c.get('sign_date'):
                    amt = f'{c["amount"]:,.0f}' if c.get('amount') else ''
                    sign_info = f'■{status}  日期:{c["sign_date"]}  金額:{amt}元'
                elif status == '預計簽約':
                    amt = f'{c["expected_amount"]:,.0f}' if c.get('expected_amount') else ''
                    sign_info = f'■預計簽約  預計日期:{c.get("expected_date","")}'
                    if amt: sign_info += f'  金額:{amt}元'
                else:
                    sign_info = f'■{status}'

                set_cell(ct.cell(ri2,0), str(ri2), size=8, bg=status_bg, align='center')
                set_cell(ct.cell(ri2,1), c.get('client',''), size=8, bg=status_bg, align='left')
                set_cell(ct.cell(ri2,2), c.get('project_name',''), size=8, bg=status_bg, align='left')
                set_cell(ct.cell(ri2,3), c.get('group_name',''), size=8, bg=status_bg, align='center')
                set_cell(ct.cell(ri2,4), sign_info, size=8, bg=status_bg, align='left')

            y_pos += tbl_h2 + Inches(0.12)

    output = io.BytesIO()
    prs.save(output); output.seek(0)
    fname = f'{year}年業務部門來自民間業務收入明細_{datetime.now().strftime("%Y%m%d")}.pptx'
    return send_file(output, as_attachment=True, download_name=fname,
                     mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation')


@app.route('/export_excel')
@login_required
def export_excel():
    year = get_current_year()
    con = get_db()
    rev_rows      = con.execute('SELECT * FROM revenue WHERE year=? ORDER BY dept, month, item', (year,)).fetchall()
    unclaim_rows  = con.execute('SELECT * FROM unclaimed WHERE year=? ORDER BY dept, month, item', (year,)).fetchall()
    contract_rows = con.execute('SELECT * FROM contracts WHERE year=? ORDER BY dept, month', (year,)).fetchall()
    con.close()

    wb = openpyxl.Workbook()
    thin   = Side(style='thin')
    border = Border(left=thin, right=thin, top=thin, bottom=thin)
    hfill  = PatternFill('solid', start_color='BDD7EE', fgColor='BDD7EE')
    hfont  = Font(name='微軟正黑體', size=10, bold=True)
    ctr    = Alignment(horizontal='center', vertical='center')

    def make_sheet(ws, headers, widths, rows_data):
        for col, (h, w) in enumerate(zip(headers, widths), 1):
            c = ws.cell(row=1, column=col, value=h)
            c.font = hfont; c.fill = hfill; c.alignment = ctr; c.border = border
            ws.column_dimensions[get_column_letter(col)].width = w
        for ri, row in enumerate(rows_data, 2):
            for ci, val in enumerate(row, 1):
                ws.cell(row=ri, column=ci, value=val).border = border

    ws1 = wb.active; ws1.title = f'{year}年收支資料'
    make_sheet(ws1, ['年度','部門','月份','項目','金額','年度目標','更新時間'],
               [8,10,8,25,14,14,20],
               [(r['year'],r['dept'],r['month'],r['item'],r['amount'],r['goal'],r['updated_at']) for r in rev_rows])

    ws2 = wb.create_sheet(f'{year}年未核銷費用')
    make_sheet(ws2, ['年度','部門','月份','項目','金額','備註','更新時間'],
               [8,10,8,20,14,20,20],
               [(r['year'],r['dept'],r['month'],r['item'],r['amount'],r['note'] or '',r['updated_at']) for r in unclaim_rows])

    ws3 = wb.create_sheet(f'{year}年合約追蹤')
    make_sheet(ws3, ['年度','部門','月份','客戶/計畫','合約金額','簽約日期','預計完成','狀態','實收金額','備註','延續下月'],
               [8,10,6,25,12,12,12,12,12,20,8],
               [(r['year'],r['dept'],r['month'],r['client'],r['amount'],r['sign_date'],r['due_date'],
                 r['status'],r['actual_amount'],r['note'] or '','是' if r['carry_next'] else '') for r in contract_rows])

    output = io.BytesIO()
    wb.save(output); output.seek(0)
    return send_file(output, as_attachment=True,
                     download_name=f'{year}年業務收支資料_{datetime.now().strftime("%Y%m%d")}.xlsx',
                     mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')

DEPLOY_TOKEN = os.environ.get('DEPLOY_TOKEN', '')

@app.route('/deploy', methods=['POST'])
def deploy():
    token = request.json.get('token', '') if request.is_json else request.form.get('token', '')
    if not DEPLOY_TOKEN or token != DEPLOY_TOKEN:
        return jsonify({'error': 'unauthorized'}), 403
    import subprocess
    try:
        pull = subprocess.run(['git', 'pull'], capture_output=True, text=True,
                              cwd=os.path.dirname(os.path.abspath(__file__)))
        wsgi = '/var/www/ttri_pythonanywhere_com_wsgi.py'
        if os.path.exists(wsgi):
            import pathlib; pathlib.Path(wsgi).touch()
        return jsonify({'ok': True, 'output': pull.stdout + pull.stderr})
    except Exception as e:
        return jsonify({'error': str(e)}), 500

init_db()  # 無論何種執行方式都先初始化 DB

@app.errorhandler(500)
def internal_error(e):
    import traceback
    tb = traceback.format_exc()
    return f'<h2>Internal Server Error</h2><pre style="font-size:13px">{tb}</pre>', 500

if __name__ == '__main__':
    import socket
    try:
        local_ip = socket.gethostbyname(socket.gethostname())
    except:
        local_ip = '127.0.0.1'
    port = int(os.environ.get('PORT', 5001))
    current_year = CURRENT_ROC_YEAR
    print('=' * 55)
    print(f'  紡織所業務部門來自民間業務收入明細（多年度版）')
    print('=' * 55)
    print(f'  本機網址 : http://127.0.0.1:{port}')
    print(f'  區域網路 : http://{local_ip}:{port}')
    print(f'  當前年度 : 民國 {current_year} 年')
    print(f'  預設帳號 : admin / admin1234')
    print('  (請登入後立即修改密碼)')
    print('=' * 55)
    app.run(host='0.0.0.0', port=port, debug=True)
